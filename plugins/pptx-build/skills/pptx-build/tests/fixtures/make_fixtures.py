#!/usr/bin/env python3
"""Regenerate the binary fixtures used by run_tests.sh.

The fixtures are committed, not generated at test time — a control that
evaporates with the session cannot catch a regression next release. This script
exists so the fixtures can be rebuilt and so their content is reviewable as code.

    python3 make_fixtures.py        # rewrites messy-deck.pptx and clean-deck.pptx

messy-deck.pptx is the deck this skill exists to fix: everything on the "Blank"
layout as free textboxes, a colored band behind each title, a grouped drawing of
boxes and arrows, a 3-D pie, hand-typed "・" bullets, topic titles, no sources.
clean-deck.pptx is build_deck.py's own output — the lossless round-trip control.
"""
import os
import subprocess
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.abspath(os.path.join(HERE, "..", "..", "assets"))


def _box(slide, x, y, w, h, lines, size=18, bold=False, color=0x1A1A1A):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(color >> 16, (color >> 8) & 0xFF, color & 0xFF)
    return tb


def _band(slide):
    """The drifting colored band — the tell this skill refuses to draw."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.28),
                                Inches(13.333), Inches(0.9))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x1F, 0x3B, 0x73)
    sh.line.fill.background()
    return sh


def messy(path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. title slide, hand-floated
    s = prs.slides.add_slide(blank)
    _box(s, 1.0, 2.6, 11, 1.4, ["2026年度 事業計画"], size=40, bold=True)
    _box(s, 1.0, 4.1, 11, 0.6, ["経営企画部 2026-04-01"], size=18, color=0x6B7280)

    # 2. bullets under a band, with hand-typed glyphs and a topic title
    s = prs.slides.add_slide(blank)
    _band(s)
    _box(s, 0.9, 0.42, 11.5, 0.7, ["市場環境"], size=28, bold=True, color=0xFFFFFF)
    _box(s, 0.9, 1.6, 11.5, 4.2, [
        "・国内市場は2年連続で縮小している",
        "・海外は年率12%で成長",
        "・競合A社が価格を引き下げた",
        "・当社シェアは18%で横ばい",
    ], size=20)
    _box(s, 0.9, 6.6, 6, 0.4, ["出典: 業界統計 2026年版"], size=11, color=0x6B7280)
    _box(s, 12.6, 6.9, 0.5, 0.35, ["2"], size=11, color=0x6B7280)

    # 3. a grouped drawing of boxes and arrows — structure that cannot survive
    s = prs.slides.add_slide(blank)
    _band(s)
    _box(s, 0.9, 0.42, 11.5, 0.7, ["提供プロセス"], size=28, bold=True, color=0xFFFFFF)
    shapes = []
    for i, label in enumerate(["受注", "設計", "製造", "納品"]):
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1.0 + i * 3.0), Inches(3.0), Inches(2.4), Inches(1.2))
        b.text_frame.text = label
        b.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
        shapes.append(b)
        if i < 3:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.45 + i * 3.0),
                                   Inches(3.4), Inches(0.5), Inches(0.4))
            shapes.append(a)
    s.shapes._spTree  # noqa — grouping below needs the tree materialized
    group = s.shapes.add_group_shape(shapes)
    group.name = "process-diagram"

    # 4. a 3-D pie with too many slices
    s = prs.slides.add_slide(blank)
    _band(s)
    _box(s, 0.9, 0.42, 11.5, 0.7, ["売上構成"], size=28, bold=True, color=0xFFFFFF)
    data = CategoryChartData()
    data.categories = ["製品A", "製品B", "製品C", "製品D", "製品E", "製品F", "製品G"]
    data.add_series("2026年度", (32.0, 24.0, 15.0, 11.0, 8.0, 6.0, 4.0))
    s.shapes.add_chart(XL_CHART_TYPE.PIE_EXPLODED, Inches(3.0), Inches(1.6),
                       Inches(7.0), Inches(5.0), data)

    prs.save(path)
    print("wrote %s" % path)


def clean(path):
    spec = os.path.join(HERE, "clean-deck.yaml")
    subprocess.check_call([sys.executable, os.path.join(ASSETS, "build_deck.py"),
                           spec, "-o", path], stdout=subprocess.DEVNULL)
    print("wrote %s" % path)


if __name__ == "__main__":
    messy(os.path.join(HERE, "messy-deck.pptx"))
    clean(os.path.join(HERE, "clean-deck.pptx"))
