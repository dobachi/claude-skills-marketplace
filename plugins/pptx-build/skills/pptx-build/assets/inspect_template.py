#!/usr/bin/env python3
"""inspect_template.py — show a template's layouts and placeholders, then emit a
suggested role->placeholder map you can hand to build_deck.py --map.

This is the tool to run FIRST when "the template's title placeholder isn't being
filled": it tells you exactly which layout index and which placeholder idx carry
the title / body / subtitle / picture in *this* template, so the map is grounded
in reality instead of guessed.

Usage:
  python3 inspect_template.py corp.pptx                 # human-readable dump
  python3 inspect_template.py corp.pptx --map > map.json # write a starter map
"""
import argparse
import json
import sys

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


def ph_type(ph):
    try:
        return str(ph.placeholder_format.type).split()[0].split(".")[-1]
    except Exception:
        return "?"


def emu_in(v):
    return round((v or 0) / 914400.0, 2)


def dump(prs):
    print("=== %d slide layouts (master: %s) ===" %
          (len(prs.slide_layouts), prs.slide_masters[0].name or "?"))
    for i, lo in enumerate(prs.slide_layouts):
        print("\n[%d] layout %r" % (i, lo.name))
        phs = list(lo.placeholders)
        if not phs:
            print("     (no placeholders)")
        for ph in phs:
            pf = ph.placeholder_format
            print("     idx=%-2d type=%-13s name=%-28r  pos=(%.2f,%.2f) size=(%.2f x %.2f)\"" % (
                pf.idx, ph_type(ph), ph.name,
                emu_in(ph.left), emu_in(ph.top), emu_in(ph.width), emu_in(ph.height)))


def suggest_map(prs):
    """Build a starter map: for each slide type, pick a layout and the idx of the
    placeholders that best match each role. Edit before use if a guess is wrong."""
    layouts = list(prs.slide_layouts)

    def idx_of(lo, *types):
        for ph in lo.placeholders:
            if ph.placeholder_format.type in types:
                return ph.placeholder_format.idx
        return None

    def find_layout(*words):
        for i, lo in enumerate(layouts):
            if any(w in (lo.name or "").lower() for w in words):
                return i
        return None

    def body_idxs(lo):
        out = []
        for ph in lo.placeholders:
            if ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                out.append(ph.placeholder_format.idx)
        return out

    def find_layout_with(*types):
        for i, lo in enumerate(layouts):
            if any(ph.placeholder_format.type in types for ph in lo.placeholders):
                return i
        return None

    content_i = find_layout("title and content", "content") or 1
    content = layouts[content_i]
    title_i = find_layout("title slide") or 0
    section_i = find_layout("section") or content_i
    two_i = find_layout("two content", "comparison", "2 content") or content_i
    img_i = find_layout_with(PP_PLACEHOLDER.PICTURE) or find_layout("picture", "caption") or content_i

    cbody = body_idxs(content)
    tbody = body_idxs(layouts[two_i])
    m = {
        "title": {"layout": title_i,
                  "title": idx_of(layouts[title_i], PP_PLACEHOLDER.CENTER_TITLE,
                                  PP_PLACEHOLDER.TITLE),
                  "subtitle": idx_of(layouts[title_i], PP_PLACEHOLDER.SUBTITLE)},
        "section": {"layout": section_i,
                    "title": idx_of(layouts[section_i], PP_PLACEHOLDER.TITLE,
                                    PP_PLACEHOLDER.CENTER_TITLE)},
        "bullets": {"layout": content_i,
                    "title": idx_of(content, PP_PLACEHOLDER.TITLE),
                    "body": cbody[0] if cbody else None},
        "two_col": {"layout": two_i,
                    "title": idx_of(layouts[two_i], PP_PLACEHOLDER.TITLE),
                    "left": tbody[0] if len(tbody) > 0 else None,
                    "right": tbody[1] if len(tbody) > 1 else None},
        "big_number": {"layout": content_i,
                       "title": idx_of(content, PP_PLACEHOLDER.TITLE),
                       "body": cbody[0] if cbody else None},
        "quote": {"layout": content_i,
                  "title": idx_of(content, PP_PLACEHOLDER.TITLE),
                  "body": cbody[0] if cbody else None},
        "image": {"layout": img_i,
                  "title": idx_of(layouts[img_i], PP_PLACEHOLDER.TITLE),
                  "image": idx_of(layouts[img_i], PP_PLACEHOLDER.PICTURE),
                  "caption": idx_of(layouts[img_i], PP_PLACEHOLDER.BODY)},
        "blank": {"layout": find_layout("blank") or content_i},
    }
    # drop None-valued role keys so the map only pins what was actually found
    return {t: {k: v for k, v in roles.items() if v is not None} for t, roles in m.items()}


def main(argv=None):
    ap = argparse.ArgumentParser(description="Inspect a .pptx/.potx template.")
    ap.add_argument("template")
    ap.add_argument("--map", action="store_true", help="emit a starter map JSON to stdout")
    a = ap.parse_args(argv)
    prs = Presentation(a.template)
    if a.map:
        json.dump(suggest_map(prs), sys.stdout, ensure_ascii=False, indent=2)
        sys.stdout.write("\n")
    else:
        dump(prs)
        print("\nRun with --map to print a starter map.json for build_deck.py --map")


if __name__ == "__main__":
    main()
