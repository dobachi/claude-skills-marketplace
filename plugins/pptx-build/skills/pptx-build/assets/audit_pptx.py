#!/usr/bin/env python3
"""audit_pptx.py — audit a FINISHED .pptx for master/layout hygiene.

`validate_deck.py` lints the spec. This lints the artifact, and catches the one
failure that spec-linting cannot see: a deck whose text was floated onto blank
slides as free textboxes instead of written into the layouts' placeholders —
what you get when a deck is assembled by hand-written python-pptx instead of by
`build_deck.py`, or when template-fill lands on the wrong layout.

    python3 audit_pptx.py deck.pptx
    python3 audit_pptx.py deck.pptx --quiet     # findings only, no per-slide table

Exit code is 1 when any ERROR is found, so it can gate a build.

What it checks
  ERROR  a slide carries body text but NO placeholder holds any of it (composed
         parts named `part/...` are structure, not floated text — see below)
  ERROR  a slide is built on a layout that has no placeholders at all, yet has text
  WARN   free textboxes in the content area alongside filled placeholders
  WARN   a placeholder overrides its layout geometry (slide-level <a:xfrm>)
  WARN   an empty placeholder is left on the slide ("Click to add text")
  WARN   slides drawn from more than one slide master
  WARN   a saturated fill covers too much of a slide (the accent is emphasis,
         not a surface) — the dark page's own background is excluded
  INFO   a run of consecutive slides with the SAME skeleton (monotony is a
         design defect the same way clutter is)
  INFO   per-slide layout / placeholder / free-shape counts
"""
import argparse
import sys

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

ERROR, WARN, INFO = "ERROR", "WARN", "INFO"
_RANK = {ERROR: 0, WARN: 1, INFO: 2}

# Text this low on the page is a footnote (source line, page number) — a free
# textbox there is an annotation, not the slide's content.
FOOTER_BAND = 0.86        # fraction of slide height
DECORATION_CHARS = 6      # a free box this short is a marker, not body content
# build_deck names every composed part `part/<kind>` (card, step, quadrant,
# arrow). Those are a graphic built FROM the body placeholder's region, not text
# someone floated onto the slide — the distinction this auditor exists to make.
PART_PREFIX = "part/"
# A fill this saturated is the accent, not a neutral surface.
SATURATION = 24
# Above this fraction of the slide, a shape IS the background (the dark page).
BACKDROP = 0.90
# Accent fills past this share of a slide stop being emphasis.
ACCENT_AREA = 0.10
# Runs of identical skeletons: this many is a rhythm worth reporting.
SAME_RUN_INFO, SAME_RUN_WARN = 4, 6


def _has_ph_marker(shape):
    """True for a table/chart graphicFrame that carries a <p:ph> placeholder marker."""
    el = shape._element
    nv = el.find(qn("p:nvGraphicFramePr"))
    if nv is None:
        return False
    nv_pr = nv.find(qn("p:nvPr"))
    return nv_pr is not None and nv_pr.find(qn("p:ph")) is not None


def _slide_level_geometry(ph):
    """True when the slide's own copy of a placeholder pins its position/size,
    overriding the layout (PowerPoint: the shape no longer follows the master)."""
    sp_pr = ph._element.find(qn("p:spPr"))
    return sp_pr is not None and sp_pr.find(qn("a:xfrm")) is not None


def _fill_rgb(shape):
    """The shape's solid fill as (r, g, b), or None."""
    try:
        if shape.fill.type != 1:            # MSO_FILL.SOLID
            return None
        rgb = shape.fill.fore_color.rgb
    except Exception:
        return None
    return (rgb[0], rgb[1], rgb[2]) if rgb is not None else None


def _skeleton(layout_name, kinds, has_gfx, has_pic):
    """What a slide looks like from across the room: its layout, the parts drawn
    on it, and whether it carries a figure. Slides with the same skeleton read as
    the same slide — fine for like content, flattening when it never changes."""
    return (layout_name, tuple(sorted(kinds)), has_gfx, has_pic)


def audit(path):
    prs = Presentation(path)
    findings = []
    rows = []
    masters = set()
    page_h = prs.slide_height or Emu(6858000)
    page_w = prs.slide_width or Emu(12192000)
    page_area = float(page_h) * float(page_w)
    skeletons = []

    def add(sev, idx, msg):
        findings.append((sev, idx, msg))

    for i, slide in enumerate(prs.slides, start=1):
        layout = slide.slide_layout
        masters.add(layout.slide_master.name or id(layout.slide_master))
        layout_phs = len(list(layout.placeholders))

        ph_text, ph_empty, ph_pinned, free_body, free_foot, graphics = 0, 0, 0, 0, 0, 0
        parts = 0
        accent_area, kinds, has_pic = 0.0, set(), False
        for sh in slide.shapes:
            named_part = (sh.name or "").startswith(PART_PREFIX)
            if named_part:
                kinds.add((sh.name or "")[len(PART_PREFIX):])
            area = float(sh.width or 0) * float(sh.height or 0) / (page_area or 1)
            rgb = _fill_rgb(sh)
            if rgb and max(rgb) - min(rgb) >= SATURATION and area < BACKDROP:
                accent_area += area
            if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
                has_pic = True
            if sh.is_placeholder or _has_ph_marker(sh):
                if not sh.has_text_frame:
                    graphics += 1                     # table/chart in a placeholder
                elif sh.text_frame.text.strip():
                    ph_text += 1
                else:
                    ph_empty += 1
                # A placeholder an archetype deliberately re-placed (the narrow
                # column of a `split`) carries the part name: that geometry is the
                # composition, not a drifted box.
                if sh.is_placeholder and not named_part and _slide_level_geometry(sh):
                    ph_pinned += 1
                continue
            if named_part:
                parts += 1
                continue
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue                              # hairline, picture, spacer
            text = sh.text_frame.text.strip()
            in_footer = (sh.top or 0) >= page_h * FOOTER_BAND
            if in_footer or len(text) <= DECORATION_CHARS:
                free_foot += 1
            else:
                free_body += 1

        rows.append((i, layout.name, ph_text, graphics, parts, ph_empty, free_body, free_foot))
        skeletons.append(_skeleton(layout.name, kinds, graphics > 0, has_pic))
        if accent_area > ACCENT_AREA:
            add(WARN, i, "saturated fill covers %d%% of the slide — the accent is "
                         "emphasis, not a surface; tint it or shrink it"
                % round(accent_area * 100))

        if parts and not ph_text:
            add(WARN, i, "%d composed part(s) but no text in any placeholder — the "
                         "slide's title should be in the title placeholder" % parts)
        if free_body and not (ph_text or graphics or parts):
            add(ERROR, i, "content is in free textboxes, not placeholders (layout %r) — "
                          "build this slide from a layout instead of floating text boxes"
                % layout.name)
        elif free_body:
            add(WARN, i, "%d free textbox(es) in the content area alongside placeholders — "
                         "move that text into the layout's placeholders" % free_body)
        if layout_phs == 0 and (free_body or free_foot):
            add(ERROR, i, "layout %r has no placeholders at all — the deck is not "
                          "master-governed on this slide" % layout.name)
        if ph_pinned:
            add(WARN, i, "%d placeholder(s) pin their own position/size, overriding the "
                         "layout — edits to the layout will not move them" % ph_pinned)
        if ph_empty:
            add(WARN, i, "%d empty placeholder(s) left on the slide — fill or delete them"
                % ph_empty)

    # Rhythm: consecutive slides that look identical from across the room.
    run_start, run_len = 1, 1
    def flush(end):
        if run_len >= SAME_RUN_INFO:
            add(WARN if run_len >= SAME_RUN_WARN else INFO, run_start,
                "%d consecutive slides share the same skeleton (%s) — vary the "
                "composition where the content's shape differs, or mark the turns "
                "with a section/statement" % (run_len, skeletons[run_start - 1][0]))
    for n in range(1, len(skeletons)):
        if skeletons[n] == skeletons[n - 1]:
            run_len += 1
        else:
            flush(n)
            run_start, run_len = n + 1, 1
    flush(len(skeletons))

    if len(masters) > 1:
        add(WARN, 0, "slides come from %d different slide masters — decks normally use one"
            % len(masters))
    if not rows:
        add(ERROR, 0, "the file has no slides")
    return rows, findings


def report(path, rows, findings, quiet=False):
    if not quiet:
        print("Slide inventory — %s" % path)
        print("  %-4s %-24s %5s %5s %5s %5s %5s %5s" %
              ("#", "layout", "PH✓", "GFX", "PART", "PH∅", "FREE", "foot"))
        for i, name, ph_text, gfx, parts, ph_empty, free_body, free_foot in rows:
            print("  %-4d %-24s %5d %5d %5d %5d %5d %5d"
                  % (i, (name or "?")[:24], ph_text, gfx, parts, ph_empty,
                     free_body, free_foot))
        print("  PH✓ text in placeholders · GFX table/chart in a placeholder · "
              "PART composed part (card/step/quadrant)\n  PH∅ empty placeholder · "
              "FREE free textbox in the content area · foot footnote-band textbox\n")

    order = sorted(findings, key=lambda f: (_RANK[f[0]], f[1]))
    if order:
        print("Findings")
        for sev, idx, msg in order:
            where = "slide %-3d" % idx if idx else "deck     "
            print("  %-5s %s %s" % (sev, where, msg))
    else:
        print("Findings: none — every slide is built on a layout and writes into its "
              "placeholders.")
    n_err = sum(1 for f in findings if f[0] == ERROR)
    n_warn = sum(1 for f in findings if f[0] == WARN)
    print("\nSummary: %d error, %d warn (%d slides)" % (n_err, n_warn, len(rows)))
    return n_err


def main(argv=None):
    ap = argparse.ArgumentParser(description="Audit a .pptx for master/layout hygiene.")
    ap.add_argument("pptx")
    ap.add_argument("--quiet", action="store_true", help="findings only")
    a = ap.parse_args(argv)
    rows, findings = audit(a.pptx)
    sys.exit(1 if report(a.pptx, rows, findings, a.quiet) else 0)


if __name__ == "__main__":
    main()
