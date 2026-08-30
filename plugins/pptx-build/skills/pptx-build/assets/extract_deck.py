#!/usr/bin/env python3
"""extract_deck.py — read an EXISTING .pptx back into a build_deck spec.

This is the entry point for *refactoring* a deck instead of authoring one. Given
any .pptx — one you received, one PowerPoint produced, one an AI floated onto
blank slides — it recovers the CONTENT (titles, bullets, tables, charts, images,
notes) as a spec YAML, so the deck can be re-argued in text and re-rendered by
`build_deck.py` into master-governed, placeholder-based slides.

    python3 extract_deck.py old.pptx -o deck.yaml
    python3 extract_deck.py old.pptx -o deck.yaml --media-dir figures
    python3 extract_deck.py old.pptx -o deck.yaml --slides 3-9,12

Why extract instead of patching in place: patching keeps whatever the original
did wrong (free textboxes, drifting bands, decorative shapes, a layout nothing
follows). Extraction keeps only what carries meaning, and the rebuild puts it
back under the master. What the spec cannot express is not silently dropped —
every loss is reported, with the slide number, so you decide what to redraw.

Exit codes
  0  every slide mapped onto a spec type with no reported loss
  1  extracted, but content could not be represented (LOSS findings) — read them
     before rebuilding: SmartArt, grouped drawings, WordArt, media, unknown charts
  2  the file cannot be read as a .pptx

The report is written to stderr, the spec to -o (or stdout with `-o -`), so
`extract_deck.py in.pptx -o - > deck.yaml` stays clean.
"""
import argparse
import os
import re
import sys

try:
    import yaml
except ImportError:                                    # pragma: no cover
    sys.stderr.write("PyYAML is required: pip install -r requirements.txt\n")
    sys.exit(2)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

LOSS, WARN, INFO = "LOSS", "WARN", "INFO"
_RANK = {LOSS: 0, WARN: 1, INFO: 2}

# A shape whose top edge is below this fraction of the page is furniture:
# source line, page number, confidentiality footer.
FOOTER_BAND = 0.84
# Text this short in the content area is a label/marker, not a message.
MARKER_CHARS = 3
# A picture smaller than this fraction of the slide area is an icon or a logo,
# not the slide's figure — icons are decoration and are not carried over.
FIGURE_AREA = 0.06
# Font size (pt) at or above which a short numeric string reads as a big number.
BIG_NUMBER_PT = 32

# accent1 of the stock Office themes — finding one means the deck inherited the
# default palette, so it is not a brand color worth carrying into the rebuild.
STOCK_ACCENTS = {"4F81BD", "5B9BD5", "4472C4", "156082"}

SOURCE_RE = re.compile(r"^\s*(出典|出所|参考|注\d*|※|source|sources|note)\s*[:：]?\s*", re.I)
PAGENUM_RE = re.compile(r"^\s*[-–—]?\s*\d{1,3}\s*(/\s*\d{1,3})?\s*[-–—]?\s*$")
NUMBERISH_RE = re.compile(r"^[\d,.\s]+\s*(%|％|倍|億|万|兆|件|人|社|名|pt|ppt|x|×|円|ドル|\$|¥)?$", re.I)
SECTION_NUM_RE = re.compile(r"^\s*(\d{1,2}|[IVX]{1,4})[.．、)]?\s*$")
QUOTE_OPEN = "\"'“”„«「『"


# ---------------------------------------------------------------------------
# Findings
# ---------------------------------------------------------------------------
class Report(object):
    def __init__(self):
        self.items = []

    def add(self, sev, idx, msg):
        self.items.append((sev, idx, msg))

    @property
    def n_loss(self):
        return sum(1 for s, _, _ in self.items if s == LOSS)


# ---------------------------------------------------------------------------
# Shape inventory
# ---------------------------------------------------------------------------
def _in(v):
    return Emu(v or 0).inches


def _max_pt(tf):
    """Largest explicit run size in a text frame, in points (None when inherited)."""
    best = None
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                pt = r.font.size.pt
                best = pt if best is None else max(best, pt)
    return best


# A literal bullet glyph typed into the text: PowerPoint's own bullets live in
# the paragraph properties, so a glyph in the string is someone's hand-typed one.
# A pictorial glyph needs no space after it ("・国内市場は…"); an ASCII dash or
# star does, so "-5%の減少" keeps its minus sign.
GLYPH_RE = re.compile(r"^[\s\u3000]*(?:[•・●○◇◆■□▪▶‣][\s\u3000]*|[-–—*·]+[\s\u3000]+)(?=\S)")


def _paras(tf):
    """[(level, text)] — empty paragraphs dropped, hand-typed bullet glyphs stripped."""
    out = []
    for p in tf.paragraphs:
        text = "".join(r.text for r in p.runs) or p.text or ""
        text = GLYPH_RE.sub("", text.replace("\x0b", " ")).strip()
        if text:
            out.append((min(int(p.level or 0), 4), text))
    return out


def _first_bold(tf):
    """True when the first non-empty paragraph is entirely bold — the cue a column
    heading gives off in decks that have one (build_deck writes headings bold)."""
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text.strip()]
        if runs:
            return all(r.font.bold for r in runs)
    return False


def _is_smartart(shape):
    xml = shape._element.xml
    return "graphicData" in xml and "diagram" in xml


def _records(slide, page_w, page_h, idx, report):
    """Flatten a slide into content records, ordered top-to-bottom, left-to-right.

    Groups are flattened (a group is a drawing; its parts are recovered as text
    and reported as a loss, because the arrangement itself does not survive)."""
    recs = []

    def walk(shapes, in_group):
        for sh in shapes:
            st = sh.shape_type
            if st == MSO_SHAPE_TYPE.GROUP:
                report.add(LOSS, idx, "a grouped drawing was flattened — its text is "
                                      "kept, its arrangement is not; redraw it as a "
                                      "figure (document-figures) if it carried structure")
                walk(sh.shapes, True)
                continue
            if _is_smartart(sh):
                report.add(LOSS, idx, "SmartArt cannot be represented in the spec — its "
                                      "text is kept as bullets; if it encoded a real "
                                      "structure, redraw it and place it as an image")
                if sh.has_text_frame:
                    recs.append(_text_rec(sh, in_group))
                continue
            if st == MSO_SHAPE_TYPE.MEDIA:
                report.add(LOSS, idx, "embedded audio/video is not carried over")
                continue
            if sh.has_table:
                recs.append({"kind": "table", "sh": sh, "x": _in(sh.left), "y": _in(sh.top),
                             "w": _in(sh.width), "h": _in(sh.height)})
                continue
            if sh.has_chart:
                recs.append({"kind": "chart", "sh": sh, "x": _in(sh.left), "y": _in(sh.top),
                             "w": _in(sh.width), "h": _in(sh.height)})
                continue
            if st == MSO_SHAPE_TYPE.PICTURE or st == MSO_SHAPE_TYPE.LINKED_PICTURE:
                recs.append({"kind": "picture", "sh": sh, "x": _in(sh.left), "y": _in(sh.top),
                             "w": _in(sh.width), "h": _in(sh.height)})
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                recs.append(_text_rec(sh, in_group))
                continue
            # No text and not a content object: a band, a connector, an icon.
            if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.LINE):
                recs.append({"kind": "deco", "sh": sh, "x": _in(sh.left), "y": _in(sh.top),
                             "w": _in(sh.width), "h": _in(sh.height)})

    def _text_rec(sh, in_group):
        tf = sh.text_frame
        return {"kind": "text", "sh": sh, "tf": tf, "in_group": in_group,
                "sid": sh.shape_id, "bold0": _first_bold(tf),
                "name": sh.name or "", "auto": sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE,
                "paras": _paras(tf), "pt": _max_pt(tf),
                "x": _in(sh.left), "y": _in(sh.top), "w": _in(sh.width), "h": _in(sh.height)}

    walk(slide.shapes, False)
    recs.sort(key=lambda r: (round(r["y"], 1), round(r["x"], 1)))
    for r in recs:
        # A composed part low on the page (a matrix's x-axis labels) is content,
        # not a running footer — the name settles it.
        r["footer"] = (r["y"] >= page_h * FOOTER_BAND
                       and not str(r.get("name", "")).startswith("part/"))
        r["area"] = (r["w"] * r["h"]) / (page_w * page_h) if page_w and page_h else 0
    return recs


# ---------------------------------------------------------------------------
# Content readers
# ---------------------------------------------------------------------------
def _read_table(rec, idx, report):
    tbl = rec["sh"].table
    rows = [[c.text.replace("\n", " ").strip() for c in r.cells] for r in tbl.rows]
    if not rows:
        return None
    out = {}
    if getattr(tbl, "first_row", False) and len(rows) > 1:
        out["columns"] = rows[0]
        out["rows"] = rows[1:]
    else:
        out["rows"] = rows
    if any(len(getattr(r, "cells", [])) != len(rows[0]) for r in tbl.rows):
        report.add(WARN, idx, "the table has merged cells — they were unmerged into a "
                              "plain grid; check the result")
    return out


_CHART_WORDS = (("DOUGHNUT", "doughnut"), ("PIE", "pie"), ("BAR", "bar"),
                ("COLUMN", "column"), ("LINE", "line"), ("AREA", "area"))


def _chart_kind(chart, idx, report):
    ct = chart.chart_type
    name = (getattr(ct, "name", None) or str(ct)).upper()
    if "3D" in name or "THREE_D" in name:
        report.add(WARN, idx, "a 3-D chart was flattened to 2-D — 3-D distorts the "
                              "values it is supposed to show")
    if "EXPLODED" in name:
        report.add(WARN, idx, "an exploded pie was flattened — the gaps carried no "
                              "information; consider a bar chart instead")
    for word, kind in _CHART_WORDS:
        if word in name:
            return kind
    report.add(LOSS, idx, "chart type %s is not in the spec (column/bar/line/area/pie/"
                          "doughnut) — its data is kept, pick a supported type or place "
                          "it as an image" % name)
    return "column"


def _read_chart(rec, idx, report):
    chart = rec["sh"].chart
    out = {"chart": _chart_kind(chart, idx, report)}
    try:
        out["categories"] = [str(c) for c in chart.plots[0].categories]
    except Exception:
        out["categories"] = []
        report.add(WARN, idx, "chart categories could not be read — fill `categories` by hand")
    series = []
    try:
        for se in chart.series:
            vals = [None if v is None else (round(v, 6) if isinstance(v, float) else v)
                    for v in se.values]
            series.append({"name": str(se.name or ""), "values": vals})
    except Exception:
        report.add(WARN, idx, "chart series could not be read — fill `series` by hand")
    out["series"] = series
    if not out["categories"] or not series:
        report.add(WARN, idx, "the chart came out incomplete — check `categories`/`series`")
    return out


def _save_picture(rec, idx, n, media_dir, spec_dir, report):
    """Write the embedded image out and return the path to put in the spec."""
    try:
        image = rec["sh"].image
        blob, ext = image.blob, (image.ext or "png")
    except Exception:
        report.add(LOSS, idx, "an image could not be extracted (it may be linked, not "
                              "embedded) — re-supply the file and set `image:`")
        return None
    if not os.path.isdir(media_dir):
        os.makedirs(media_dir)
    name = "slide%02d-%d.%s" % (idx, n, ext)
    path = os.path.join(media_dir, name)
    with open(path, "wb") as f:
        f.write(blob)
    return os.path.relpath(path, spec_dir or ".")


# ---------------------------------------------------------------------------
# Slide classification
# ---------------------------------------------------------------------------
def _title_of(slide, texts):
    """(text, shape id, paragraphs) — the title placeholder, else the topmost
    largest single line. The id is what excludes it from the body; comparing the
    shape object would not work, since each access returns a fresh proxy."""
    try:
        t = slide.shapes.title
    except Exception:
        t = None
    if t is not None and t.has_text_frame and t.text_frame.text.strip():
        paras = _paras(t.text_frame)
        return " ".join(x[1] for x in paras), t.shape_id, paras
    upper = [r for r in texts if not r["footer"] and r["paras"]]
    if not upper:
        return "", None, []
    ranked = sorted(upper, key=lambda r: (-(r["pt"] or 0), r["y"]))
    head = ranked[0]
    if len(head["paras"]) == 1 and len(head["paras"][0][1]) <= 120:
        return head["paras"][0][1], head["sid"], head["paras"]
    return "", None, []


PART_KINDS = (("part/card", "cards"), ("part/step", "steps"),
              ("part/quadrant", "matrix"), ("part/lead", "lead"),
              ("part/rest", "lead"))


def _named_parts(texts):
    """The archetype a build_deck-produced slide was drawn from, by shape name."""
    if any(r.get("name", "").startswith("part/lead") for r in texts):
        return "lead"          # one lead + its supporting parts
    for prefix, kind in PART_KINDS:
        if len([r for r in texts if r.get("name", "").startswith(prefix)]) >= 2:
            return kind
    return None


def _dark(slide):
    """True when the slide paints itself dark — the inverted page."""
    try:
        rgb = slide.background.fill.fore_color.rgb
    except Exception:
        return False
    if rgb is None:
        return False
    r, g, b = rgb[0], rgb[1], rgb[2]
    return (0.299 * r + 0.587 * g + 0.114 * b) < 110


def _looks_like_split(figures, body):
    """One figure on one side, the reading on the other: horizontally disjoint,
    vertically overlapping. Text UNDER a figure is a caption, not a split."""
    if len(figures) != 1 or not body:
        return False
    f = figures[0]
    if f["area"] < 0.2:
        return False
    for r in body:
        if r["name"].startswith("part/split-text"):
            return True
    for r in body:
        disjoint = (r["x"] + r["w"] <= f["x"] + 0.05) or (r["x"] >= f["x"] + f["w"] - 0.05)
        overlaps = min(r["y"] + r["h"], f["y"] + f["h"]) - max(r["y"], f["y"]) > 0.3
        if not (disjoint and overlaps):
            return False
    return True


def _looks_like_cards(body):
    """A foreign deck's card row: 3-4 autoshapes of near-equal width, tops aligned,
    each holding a short amount of text. Deliberately strict — a wrong guess here
    turns prose into labels."""
    boxes = [r for r in body if r.get("auto") and r["w"] > 0]
    if not 3 <= len(boxes) <= 4 or len(boxes) != len(body):
        return False
    tops = [r["y"] for r in boxes]
    widths = [r["w"] for r in boxes]
    if max(tops) - min(tops) > 0.25 or max(widths) > min(widths) * 1.15:
        return False
    return all(sum(len(t) for _, t in r["paras"]) <= 120 and r["paras"] for r in boxes)


def _part_items_from(recs, drop_index=False):
    """[{label, text}] from a row of part shapes: first line is the label."""
    items = []
    for r in sorted(recs, key=lambda r: (round(r["y"], 1), r["x"])):
        paras = [t for _, t in r["paras"]]
        if drop_index and paras and re.match(r"^\d{1,2}$", paras[0].strip()):
            paras = paras[1:]     # the step number is drawn, not authored
        if not paras:
            continue
        items.append({"label": paras[0], "text": " ".join(paras[1:])})
    return items


def _bullets_from(recs):
    items = []
    for r in recs:
        for level, text in r["paras"]:
            items.append({"text": text, "level": level} if level else text)
    return items


def _is_body(r, title_sid):
    return (r["kind"] == "text" and r["sid"] != title_sid and not r["footer"]
            and r["paras"])


def _classify(slide, recs, i, is_first, report):
    """(type, title, body records, title paragraphs) for one slide.

    Deliberately conservative: when the shape of a slide is ambiguous it falls
    through to `bullets`, which loses no text. Every judgement it makes is one a
    human re-reads in the spec — the extractor's job is to lose nothing, not to
    guess cleverly."""
    texts = [r for r in recs if r["kind"] == "text"]
    title, title_sid, title_paras = _title_of(slide, texts)
    body = [r for r in texts if _is_body(r, title_sid)]
    layout_name = (slide.slide_layout.name or "").lower()
    out = lambda kind, t=None, b=None: (kind, title if t is None else t,
                                        body if b is None else b, title_paras)

    # A deck this generator produced names its composed parts; that name is the
    # most reliable signal there is, so it wins over every geometric guess.
    named = _named_parts(texts)
    if named:
        return out(named)
    has = lambda k: any(r["kind"] == k for r in recs)
    figures = [r for r in recs if r["kind"] == "picture" and r["area"] >= FIGURE_AREA]
    if has("table"):
        return out("table")
    if has("chart"):
        return out("chart")
    if figures:
        # Figure BESIDE its reading is a different composition from figure ABOVE
        # its caption — and the difference is the whole point of `split`.
        if _looks_like_split(figures, body):
            return out("split")
        return out("image")

    body_paras = [pr for r in body for pr in r["paras"]]
    words = sum(len(t) for _, t in body_paras)
    lone = body[0]["paras"] if len(body) == 1 else []

    # A quote reads as a quote wherever it sits — in the title placeholder (which
    # is where build_deck puts it) or alone in the body.
    head = title_paras[0][1] if title_paras else ""
    if head and head[0] in QUOTE_OPEN and len(title) <= 300:
        return out("quote")
    if not title and lone and lone[0][1][0] in QUOTE_OPEN and words <= 300:
        return out("quote", "")

    if "title slide" in layout_name:
        return out("title")
    # A sentence alone on a title-only layout is a statement, not a deck title.
    if ("title only" in layout_name or "statement" in layout_name) and title_paras \
            and not body and len(title) <= 160:
        return out("statement")
    if is_first and len(body) <= 1 and words <= 160:
        return out("title")

    if (lone and NUMBERISH_RE.match(lone[0][1]) and len(lone[0][1]) <= 12
            and len(lone) <= 2
            and (body[0]["pt"] is None or body[0]["pt"] >= BIG_NUMBER_PT)):
        return out("big_number")

    # A divider: a title and nothing else, or a title beside a bare section number.
    if title and not body:
        return out("section")
    if title and len(lone) == 1 and SECTION_NUM_RE.match(lone[0][1]):
        return out("section")
    if "section" in layout_name and words <= 80:
        return out("section", title or (lone[0][1] if lone else ""), [])

    if _looks_like_cards(body):
        report.add(WARN, i, "a row of equal boxes was read as `cards` — check that the "
                            "items really are equivalent; if they are a sequence, make "
                            "it `steps`, and if they are a list, `bullets`")
        return out("cards")

    if len(body) == 2:
        a, b = sorted(body, key=lambda r: r["x"])
        overlap = min(a["x"] + a["w"], b["x"] + b["w"]) - max(a["x"], b["x"])
        if overlap <= 0.2 and abs(a["y"] - b["y"]) <= 0.8:
            return out("two_col")

    if not body and not title:
        return out("blank")
    return out("bullets")


def _column(rec):
    """A column: a bold first line is its heading, the rest are its bullets."""
    paras = rec["paras"]
    if rec.get("bold0") and len(paras) > 1:
        return {"heading": paras[0][1], "bullets": _items(paras[1:])}
    return {"bullets": _items(paras)}


def _items(paras):
    return [{"text": t, "level": lv} if lv else t for lv, t in paras]


def _split_footnotes(recs, idx, report):
    """Pull the source line out of the footer band; drop page numbers and furniture."""
    source = None
    for r in recs:
        if r["kind"] != "text" or not r["footer"]:
            continue
        text = " ".join(t for _, t in r["paras"]).strip()
        if not text or PAGENUM_RE.match(text):
            continue
        if SOURCE_RE.match(text):
            source = text if source is None else source + " / " + text
        else:
            report.add(INFO, idx, "footer text dropped (%r) — running footers belong to "
                                  "the master, not to a slide" % text[:40])
    return source


# ---------------------------------------------------------------------------
# Extraction
# ---------------------------------------------------------------------------
def _accent_of(prs):
    """The template's accent1, when it is a real color rather than a neutral."""
    try:
        theme = prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
        xml = theme.blob.decode("utf-8", "ignore")
    except Exception:
        return None
    m = re.search(r"<a:accent1>.*?val=\"([0-9A-Fa-f]{6})\"", xml, re.S)
    if not m:
        return None
    hexv = m.group(1).upper()
    if hexv in STOCK_ACCENTS:   # the deck never chose a color; it kept Office's
        return None
    r, g, b = (int(hexv[j:j + 2], 16) for j in (0, 2, 4))
    if max(r, g, b) - min(r, g, b) < 24:       # gray/near-neutral: not an accent
        return None
    return hexv


def _parse_range(text, n):
    if not text:
        return list(range(1, n + 1))
    picked = []
    for part in text.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, _, b = part.partition("-")
            picked.extend(range(int(a), int(b) + 1))
        else:
            picked.append(int(part))
    return [i for i in picked if 1 <= i <= n]


def extract(path, media_dir, spec_dir, slides_arg=None, keep_notes=True):
    prs = Presentation(path)
    report = Report()
    page_w, page_h = _in(prs.slide_width), _in(prs.slide_height)
    wanted = _parse_range(slides_arg, len(prs.slides))

    meta = {"aspect": "4:3" if page_w and page_h / page_w > 0.7 else "16:9"}
    accent = _accent_of(prs)
    if accent:
        meta["accent"] = accent

    out_slides = []
    for i, slide in enumerate(prs.slides, start=1):
        if i not in wanted:
            continue
        recs = _records(slide, page_w, page_h, i, report)
        kind, title, body, title_paras = _classify(slide, recs, i, i == 1, report)
        body_paras = [pr for r in body for pr in r["paras"]]
        s = {"type": kind}
        if title and kind != "quote":
            s["title"] = title

        if kind == "table":
            tables = [r for r in recs if r["kind"] == "table"]
            tbl = _read_table(tables[0], i, report)
            if tbl:
                s.update(tbl)
            if len(tables) > 1:
                report.add(LOSS, i, "%d tables on one slide — only the first was kept; "
                                    "a slide carries one message" % len(tables))
        elif kind == "chart":
            charts = [r for r in recs if r["kind"] == "chart"]
            s.update(_read_chart(charts[0], i, report))
            if len(charts) > 1:
                report.add(LOSS, i, "%d charts on one slide — only the first was kept; "
                                    "split the slide" % len(charts))
        elif kind == "split":
            figs = [r for r in recs if r["kind"] == "picture" and r["area"] >= FIGURE_AREA]
            path = _save_picture(figs[0], i, 1, media_dir, spec_dir, report)
            if path:
                s["image"] = path
            if body and body[0].get("bold0") and len(body[0]["paras"]) > 1:
                s["heading"] = body_paras[0][1]
                s["bullets"] = _items(body_paras[1:])
            else:
                s["bullets"] = _items(body_paras)
            if body and figs[0]["x"] < body[0]["x"]:
                s["flip"] = True          # the figure sits on the left
        elif kind == "image":
            figs = [r for r in recs if r["kind"] == "picture" and r["area"] >= FIGURE_AREA]
            path = _save_picture(figs[0], i, 1, media_dir, spec_dir, report)
            if path:
                s["image"] = path
            if len(figs) > 1:
                for n, extra in enumerate(figs[1:], start=2):
                    _save_picture(extra, i, n, media_dir, spec_dir, report)
                report.add(LOSS, i, "%d figures on one slide — only the first is in the "
                                    "spec; the rest are in the media directory, give each "
                                    "its own slide" % len(figs))
            if body_paras:
                s["caption"] = body_paras[0][1]
                if len(body_paras) > 1:
                    s["note"] = " ".join(t for _, t in body_paras[1:])
        elif kind == "lead":
            head = [r for r in body if r.get("name", "").startswith("part/lead")]
            rest = [r for r in body if r.get("name", "").startswith("part/rest")]
            items = _part_items_from(head or body[:1])
            if items:
                s["lead"] = items[0]
            s["rest"] = _part_items_from(rest or body[1:])
        elif kind in ("cards", "steps", "matrix"):
            key = {"cards": "cards", "steps": "steps", "matrix": "quadrants"}[kind]
            pre = {"cards": "part/card", "steps": "part/step",
                   "matrix": "part/quadrant"}[kind]
            parts = [r for r in body if r.get("name", "").startswith(pre)] or body
            s[key] = _part_items_from(parts, drop_index=(kind == "steps"))
            if kind == "matrix":
                ax = [r for r in body if r.get("name", "") == "part/axis"]
                if len(ax) == 4 and parts:
                    left = min(r["x"] for r in parts)
                    ys = sorted([r for r in ax if r["x"] < left], key=lambda r: r["y"])
                    xs = sorted([r for r in ax if r["x"] >= left], key=lambda r: r["x"])
                    if len(ys) == 2:
                        s["y_axis"] = [ys[1]["paras"][0][1], ys[0]["paras"][0][1]]
                    if len(xs) == 2:
                        s["x_axis"] = [xs[0]["paras"][0][1], xs[1]["paras"][0][1]]
        elif kind == "two_col":
            for key, r in zip(("left", "right"), sorted(body, key=lambda r: r["x"])):
                s[key] = _column(r)
        elif kind == "big_number":
            s["number"] = body_paras[0][1]
            if len(body_paras) > 1:
                s["caption"] = " ".join(t for _, t in body_paras[1:])
        elif kind == "quote":
            lines = [t for _, t in (title_paras or [])] + [t for _, t in body_paras]
            s["quote"] = lines[0].strip("".join(QUOTE_OPEN) + "\u201d\u2019\u300d\u300f\u00bb")
            if len(lines) > 1:      # the line after the quote is who said it
                s["attribution"] = lines[1].lstrip("\u2014\u2013-\u2010\u2015\u30fc\u203b ").strip()
        elif kind == "statement":
            lines = [t for _, t in (title_paras or [])]
            s.pop("title", None)
            s["text"] = lines[0] if lines else ""
            if len(lines) > 1:
                s["sub"] = " ".join(lines[1:])
        elif kind == "title":
            sub = " ".join(t for _, t in body_paras)
            if sub:
                s["subtitle"] = sub
        elif kind == "section":
            # A bare number beside the title is the section number, not content.
            nums = [t for _, t in body_paras if SECTION_NUM_RE.match(t)]
            if nums:
                s["number"] = nums[0].strip(".．、) ")
            elif SECTION_NUM_RE.match(title or "") and body_paras:
                s["number"] = title.strip(".．、) ")
                s["title"] = " ".join(t for _, t in body_paras)
        elif kind == "bullets":
            items = _items(body_paras)
            if items:
                s["bullets"] = items

        if kind in ("section", "statement") and _dark(slide):
            s["invert"] = True          # the dark page is a deliberate turn
        src = _split_footnotes(recs, i, report)
        if src:
            s["source"] = src
        if keep_notes and slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                s["notes"] = note

        deco = [r for r in recs if r["kind"] == "deco" and r["area"] >= 0.02]
        if deco:
            report.add(INFO, i, "%d decorative shape(s) (band, box, connector) dropped — "
                                "if any encoded structure, redraw it as a real figure"
                       % len(deco))
        icons = [r for r in recs if r["kind"] == "picture" and r["area"] < FIGURE_AREA]
        if icons:
            report.add(INFO, i, "%d small image(s) treated as icons/logos and dropped"
                       % len(icons))
        if any(r["kind"] == "text" and r.get("in_group") for r in recs):
            report.add(WARN, i, "text recovered from inside a grouped drawing landed in "
                                "the bullets — reread this slide's wording")
        out_slides.append(s)

    if not out_slides:
        report.add(LOSS, 0, "no slides were extracted — check --slides")
    return {"meta": meta, "slides": out_slides}, report


# ---------------------------------------------------------------------------
# Output
# ---------------------------------------------------------------------------
HEADER = """\
# Extracted from: %s
# by extract_deck.py — this is CONTENT recovered from an existing deck, not a
# finished spec. Before rebuilding:
#   1. read the extraction report (what could not be carried over),
#   2. rewrite topic titles into action titles that state the takeaway,
#   3. split any slide carrying more than one message,
#   4. run: python3 validate_deck.py %s
# Then: python3 build_deck.py %s -o refactored.pptx && python3 audit_pptx.py refactored.pptx
"""


def dump_spec(spec, src, out_path):
    body = yaml.safe_dump(spec, allow_unicode=True, sort_keys=False,
                          default_flow_style=False, width=100)
    name = os.path.basename(out_path) if out_path != "-" else "deck.yaml"
    return HEADER % (os.path.basename(src), name, name) + body


def print_report(report, spec, src, out_path):
    n = len(spec["slides"])
    types = {}
    for s in spec["slides"]:
        types[s["type"]] = types.get(s["type"], 0) + 1
    sys.stderr.write("Extracted %d slide(s) from %s -> %s\n" % (n, src, out_path))
    sys.stderr.write("  types: %s\n" % ", ".join(
        "%s×%d" % (k, v) for k, v in sorted(types.items())))
    items = sorted(report.items, key=lambda f: (_RANK[f[0]], f[1]))
    if items:
        sys.stderr.write("\nExtraction report\n")
        for sev, idx, msg in items:
            where = "slide %-3d" % idx if idx else "deck     "
            sys.stderr.write("  %-4s %s %s\n" % (sev, where, msg))
    n_loss = report.n_loss
    sys.stderr.write("\nSummary: %d loss, %d warn, %d info\n" % (
        n_loss,
        sum(1 for s, _, _ in report.items if s == WARN),
        sum(1 for s, _, _ in report.items if s == INFO)))
    if n_loss:
        sys.stderr.write("LOSS means the spec cannot express it — decide per item: "
                         "redraw as a figure, split the slide, or drop it.\n")
    return n_loss


def main(argv=None):
    ap = argparse.ArgumentParser(
        description="Extract an existing .pptx into a build_deck spec for refactoring.")
    ap.add_argument("pptx", help="the existing .pptx to read")
    ap.add_argument("-o", "--out", required=True, help="spec output path (.yaml), or - for stdout")
    ap.add_argument("--media-dir", help="where extracted images go (default: <out>_media)")
    ap.add_argument("--slides", help="only these slides, e.g. 3-9,12")
    ap.add_argument("--no-notes", action="store_true", help="do not carry speaker notes over")
    a = ap.parse_args(argv)

    if not os.path.exists(a.pptx):
        sys.stderr.write("no such file: %s\n" % a.pptx)
        return 2
    spec_dir = "." if a.out == "-" else (os.path.dirname(os.path.abspath(a.out)) or ".")
    media = a.media_dir or os.path.join(
        spec_dir, (os.path.splitext(os.path.basename(a.out))[0] if a.out != "-" else "deck")
        + "_media")
    try:
        spec, report = extract(a.pptx, media, spec_dir, a.slides, not a.no_notes)
    except Exception as e:                                   # not a .pptx, or corrupt
        sys.stderr.write("cannot read %s as a .pptx: %s\n" % (a.pptx, e))
        return 2

    text = dump_spec(spec, a.pptx, a.out)
    if a.out == "-":
        sys.stdout.write(text)
    else:
        with open(a.out, "w", encoding="utf-8") as f:
            f.write(text)
    return 1 if print_report(report, spec, a.pptx, a.out) else 0


if __name__ == "__main__":
    sys.exit(main())
