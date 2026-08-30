#!/usr/bin/env python3
"""build_deck.py — Generate a clean, white-based .pptx that does not look AI-made.

Engine: python-pptx. Two rendering paths, ONE spec format:

  1. Default (no --template): build from a blank presentation and place text on a
     single shared grid. White background, dark ink, one restrained accent, and a
     short accent hairline whose coordinates are computed once from the grid — so
     it is identical on every slide of a family and cannot drift. This is the
     "looks human-designed" path.

  2. Template-fill (--template corporate.pptx/.potx): OPEN the real template and
     reuse its slide layouts and PLACEHOLDERS. Each spec slide is mapped to a
     layout, and its title/body/etc. are written into that layout's placeholders,
     so the deck inherits the template's master, theme, fonts, and logos. This is
     the path that actually honors a provided corporate template.

Usage:
  python3 build_deck.py SPEC.yaml -o out.pptx
  python3 build_deck.py SPEC.yaml -o out.pptx --theme themes/brand-example.json
  python3 build_deck.py SPEC.yaml -o out.pptx --template corp.pptx [--map map.json]

Inspect a template first to build/verify a map:  python3 inspect_template.py corp.pptx
"""
import argparse
import colorsys
import json
import math
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_THEME = os.path.join(HERE, "themes", "minimal-white.json")


# ---------------------------------------------------------------------------
# Spec / theme / map loading
# ---------------------------------------------------------------------------
def load_spec(path):
    with open(path, encoding="utf-8") as f:
        raw = f.read()
    if path.lower().endswith((".yaml", ".yml")):
        import yaml
        return yaml.safe_load(raw)
    return json.loads(raw)


def load_theme(path):
    with open(path, encoding="utf-8") as f:
        return json.load(f)


# Readable type scale. Every size the default renderer uses is routed through
# here so a deck can never silently fall back to tiny text. `min_body`/`min_title`
# are floors the renderer will NOT shrink past — when content won't fit at the
# floor it warns to split the slide instead of shrinking (see _check_overflow).
SIZE_DEFAULTS = {
    "title_max": 34, "title_min": 24, "title_slide": 40, "subtitle": 20,
    "section": 34, "section_number": 22,
    "body": 18, "body_sub": 16, "min_body": 16,
    "big_number": 88, "big_caption": 20,
    "quote": 28, "quote_attr": 18,
    "caption": 15, "caption_note": 14,
    "table": 15, "table_header": 15, "chart_label": 12,
    "source": 11, "page_number": 11,
    # composed parts (cards / steps / matrix) and the statement archetype
    "part_label": 18, "part_text": 16, "part_index": 13, "axis": 13,
    "statement": 32, "statement_sub": 18,
}

# --- shape tokens -----------------------------------------------------------
# One spacing unit, one corner radius, one line weight — for the whole deck.
# "Tidy" is not a matter of adding parts; it is every part agreeing on these
# three numbers. 1 unit = 0.0833 in = 8 px at 96 dpi, and every gap, padding and
# offset below is an integer multiple of it.
SHAPE_DEFAULTS = {
    "unit": 0.0833,     # spacing scale base (in)
    "radius": 0.06,     # the only corner radius (in); 0 = square corners
    "line": 0.75,       # the only line weight (pt)
}


def _shape_defaults(theme, meta=None):
    sh = dict(SHAPE_DEFAULTS)
    sh.update(theme.get("shape") or {})
    for k, v in ((meta or {}).get("shape") or {}).items():
        if isinstance(v, (int, float)):
            sh[k] = v
    theme["shape"] = sh
    for k in ("surface", "surface_hi", "border", "invert_bg"):
        theme["color"].setdefault(k, "auto")
    return theme


def _u(theme, n=1):
    """n spacing units, in inches. Every gap in a composed part comes from here."""
    return theme["shape"]["unit"] * n


# --- one hue, several tones ------------------------------------------------
# Mixing a color with white drops its SATURATION as well as raising its
# lightness, so pale steps come out grey and muddy — visibly so behind thin
# Japanese strokes. These work in HSL: lightness is set to an explicit target
# and saturation is held (scaled down only enough to keep pale steps from
# looking like poster paint). Every surface, border and pale series in the deck
# is one of these, so a deck really is one hue at several weights.
TONE = {                     # target lightness for each named role
    "surface": 0.965,        # a part's ground
    "surface_hi": 0.930,     # the highlighted part's ground
    "border": 0.900,         # a part's edge
    "invert": 0.120,         # the dark page
}
# Series ramp for a `tonal` chart: absolute lightness steps, >= 0.13 apart, so
# the bars stay distinguishable in projection and in greyscale.
SERIES_TONES = (None, 0.58, 0.72, 0.85)     # None = the accent's own lightness


def _tone(hex6, lightness, sat_scale=None):
    """The same hue at an explicit lightness. `sat_scale=None` keeps a sensible
    amount of saturation for that lightness."""
    hex6 = str(hex6).lstrip("#")
    r, g, b = (int(hex6[i:i + 2], 16) / 255.0 for i in (0, 2, 4))
    h, _l, sat = colorsys.rgb_to_hls(r, g, b)
    lightness = max(0.0, min(1.0, lightness))
    if sat_scale is None:
        # Pale steps need less saturation to stay quiet; dark steps keep theirs.
        sat_scale = 0.35 + 0.65 * (1.0 - lightness)
    r, g, b = colorsys.hls_to_rgb(h, lightness, max(0.0, min(1.0, sat * sat_scale)))
    return "%02X%02X%02X" % tuple(int(round(v * 255)) for v in (r, g, b))


def _tint(hex6, pct):
    """Kept for callers that think in "how much color": 0 = paper, 1 = the color."""
    return _tone(hex6, 1.0 - 0.58 * max(0.0, min(1.0, pct)))


def _apply_size_defaults(theme, meta=None):
    """Ensure theme['size'] has every key, then let meta.size.* override per deck.
    Shape tokens (unit / radius / line) and the derived neutrals ride along —
    one call, one place."""
    _shape_defaults(theme, meta)
    _resolve_colors(theme)
    sz = dict(SIZE_DEFAULTS)
    sz.update(theme.get("size") or {})
    for k, v in ((meta or {}).get("size") or {}).items():
        if isinstance(v, (int, float)):
            sz[k] = v
    theme["size"] = sz
    return theme


def _inverted(theme):
    """The dark page. A deck of white slides has no contrast at the DECK level —
    every page weighs the same — so the turns in the argument get one inverted
    slide. This is not the drifting band: it is the whole page, on the two types
    that mark a turn (`section`, `statement`), so there is no edge to misalign.
    Every color is the same hue at a different lightness, not a new palette."""
    t = dict(theme)
    src, c = theme["color"], dict(theme["color"])
    c["bg"] = src["invert_bg"]
    c["ink"] = src["invert_ink"]
    c["muted"] = src["invert_muted"]
    c["accent"] = _tone(src["accent"], 0.66)      # the rule has to read on dark
    c["surface"] = _tone(src["accent"], 0.20)
    c["surface_hi"] = _tone(src["accent"], 0.28)
    c["border"] = _tone(src["accent"], 0.32)
    t["color"] = c
    return t


def _resolve_colors(theme):
    """Neutrals are derived from the accent unless the theme pins them.

    A pure-grey surface next to a colored accent reads as two colors; the same
    surface carried 4% toward the accent reads as one. `"auto"` (the default)
    asks for the derived value; a literal hex opts out."""
    c = theme["color"]
    for key, role in (("surface", "surface"), ("surface_hi", "surface_hi"),
                      ("border", "border"), ("invert_bg", "invert")):
        if str(c.get(key, "auto")).lower() in ("", "auto", "none"):
            c[key] = _tone(c["accent"], TONE[role])
    c.setdefault("invert_ink", "FFFFFF")
    c.setdefault("invert_muted", _tone(c["accent"], 0.72, sat_scale=0.25))
    return theme


def apply_meta(theme, meta):
    """meta.* in the spec overrides individual theme values for one deck."""
    meta = meta or {}
    for k in ("bg", "ink", "muted", "accent"):
        if meta.get(k):
            theme["color"][k] = str(meta[k]).lstrip("#")
    for k, dst in (("font_heading", "heading"), ("font_body", "body"), ("font_number", "number")):
        if meta.get(k):
            theme["font"][dst] = meta[k]
    if isinstance(meta.get("rule"), bool):
        theme["rule"] = meta["rule"]
    if isinstance(meta.get("page_numbers"), bool):
        theme["pageNumbers"] = meta["page_numbers"]
    if meta.get("aspect") in ("16:9", "4:3"):
        theme["aspect"] = meta["aspect"]
    return theme


# ---------------------------------------------------------------------------
# Grid — single source of truth, in inches.
# ---------------------------------------------------------------------------
def make_grid(theme):
    wide = theme.get("aspect", "16:9") != "4:3"
    over = theme.get("grid") or {}
    g = {
        "pageW": 13.33 if wide else 10.0,
        "pageH": 7.5,
        "marginX": over.get("marginX", 0.92 if wide else 0.75),
        "top": over.get("top", 0.62),
        # titleH is the fixed title REGION height — tall enough for up to two
        # lines at title_max. The title text is BOTTOM-anchored inside it and
        # shares one baseline across every slide, so a two-line title grows
        # UPWARD into the top margin and can never reach the hairline or body.
        "titleH": over.get("titleH", 1.04),
        "gap": 0.22,
        "footY": 7.04,
        "ruleLen": 1.05,
        "ruleH": 0.045,
    }
    g["contentW"] = g["pageW"] - 2 * g["marginX"]
    g["titleBottom"] = g["top"] + g["titleH"]
    # Hairline sits in the gap below the title region; body starts below the hairline.
    g["bodyTop"] = g["titleBottom"] + g["gap"]
    g["bodyH"] = g["footY"] - g["bodyTop"] - 0.1
    g["ruleY"] = {
        "content": g["titleBottom"] + 0.08,
        "title": 2.30,
        "section": 2.78,
        "quote": 2.02,
    }
    return g


# ---------------------------------------------------------------------------
# Overflow guard — keep type LARGE. We never auto-shrink below the floor; when
# content won't fit at readable sizes we warn so the operator splits the slide.
# ---------------------------------------------------------------------------
# Directory of the spec file. A spec that came out of extract_deck.py points at
# images next to itself (`deck_media/slide03-1.png`), so a relative path must
# resolve against the spec, not only against wherever the build was run from.
_BASE_DIR = None


def _asset(path):
    """Resolve a spec-relative asset path: as given first, then next to the spec."""
    if not path or os.path.isabs(path) or os.path.exists(path) or not _BASE_DIR:
        return path
    near = os.path.join(_BASE_DIR, path)
    return near if os.path.exists(near) else path


def _apply_notes(prs, slides):
    """Carry `notes:` into the real notes slide — what the presenter says is part
    of the deck, and a refactor that drops it loses half the argument."""
    for slide, s in zip(prs.slides, slides):
        text = (s or {}).get("notes")
        if text:
            slide.notes_slide.notes_text_frame.text = str(text)


_WARNINGS = []


def _warn(msg):
    _WARNINGS.append(msg)
    sys.stderr.write("  warning: %s\n" % msg)


def _disp_width(text):
    """Approximate display width of text in 'em' units: CJK/full-width ~1.0, ASCII ~0.55."""
    w = 0.0
    for ch in str(text):
        w += 1.0 if ord(ch) > 0x2E7F else 0.55
    return w


def _est_lines(text, size_pt, width_in):
    """Estimate how many wrapped lines `text` takes at `size_pt` in `width_in`."""
    if not text:
        return 0
    chars_per_line = max(1.0, width_in / (size_pt / 72.0))
    n = 0
    for seg in str(text).split("\n"):
        n += max(1, math.ceil(_disp_width(seg) / chars_per_line))
    return n


# ---------------------------------------------------------------------------
# Low-level text helpers (the EA-font handling is what keeps Japanese on-brand).
# ---------------------------------------------------------------------------
def _set_run_font(run, name=None, size=None, bold=None, color=None):
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = RGBColor.from_string(color)
    if name:
        f.name = name  # inserts <a:latin> in the correct schema position
        rPr = run._r.get_or_add_rPr()
        latin = rPr.find(qn("a:latin"))
        for tag in ("a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {})
                latin.addnext(el)  # ea/cs follow latin per the schema
            el.set("typeface", name)


def _bullet_para(p, marker, level, font_name, space_after=9):
    """Give a paragraph a real hanging-indent bullet (no glyph flush, no tofu)."""
    p.level = level
    p.space_after = Pt(space_after)
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(0.30 + 0.22 * level))))
    pPr.set("indent", str(-int(Inches(0.22))))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": font_name}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": marker}))


def _no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buFont"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    if pPr.find(qn("a:buNone")) is None:
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, name=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        box.name = name            # `part/...` marks it as composed structure
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def set_simple(tf, text, theme, font="body", size=18, bold=False, color="ink",
               align=PP_ALIGN.LEFT, line_spacing=None):
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    _no_bullet(p)
    run = p.add_run()
    run.text = text
    _set_run_font(run, name=theme["font"][font], size=size, bold=bold,
                  color=theme["color"][color])
    return p


# ---------------------------------------------------------------------------
# DEFAULT MODE — build on the STANDARD PowerPoint layout set and write into each
# layout's real placeholders (title / body / subtitle / picture). We never
# free-float a textbox onto a blank slide: every title AND every body lives in a
# master-governed placeholder.
#
# The placeholder GEOMETRY is configured ONCE per layout (setup_layouts). Every
# slide created from a layout inherits that geometry, so both titles and body
# content are master-governed: reposition a placeholder on the layout later in
# PowerPoint and every slide built on it moves together. Title regions are
# bottom-anchored, so a long (two-line) title grows UP into the top margin and
# can never collide with the hairline or the body. Explicit run sizes (with the
# readable floors) and autofit-off keep type large — the template's own
# shrink-to-fit can never make our text tiny.
# ---------------------------------------------------------------------------
TITLE_LAYOUT = 0        # "Title Slide"          — CENTER_TITLE + SUBTITLE
CONTENT_LAYOUT = 1      # "Title and Content"    — TITLE + body (OBJECT)
SECTION_LAYOUT = 2      # "Section Header"       — TITLE + body (BODY)
TWO_CONTENT_LAYOUT = 3  # "Two Content"          — TITLE + two bodies (OBJECT)
IMAGE_LAYOUT = 8        # "Picture with Caption" — TITLE + PICTURE + caption (BODY)
QUOTE_LAYOUT = 5        # "Title Only"           — its TITLE placeholder holds the quote
BLANK_LAYOUT = 6        # "Blank"                — only for type: blank


def _normalize_bullets(items):
    out = []
    for it in items or []:
        if isinstance(it, str):
            out.append({"text": it, "level": 0})
        else:
            out.append({"text": it.get("text", ""), "level": it.get("level", 0)})
    return out


def _place(ph, x, y, w, h):
    ph.left, ph.top, ph.width, ph.height = Inches(x), Inches(y), Inches(w), Inches(h)


def _phs_by_role(container):
    """(title_ph, [body_ph, ...]) for a slide OR a layout. Bodies are the
    content/text/subtitle placeholders in reading order; footer/date/slide-number
    placeholders are ignored. Works on both because both expose `.placeholders`."""
    title, bodies = None, []
    for ph in container.placeholders:
        t = ph.placeholder_format.type
        if t in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            title = ph
        elif t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE):
            bodies.append(ph)
    bodies.sort(key=lambda p: (p.top or 0, p.left or 0))
    return title, bodies


def _picture_ph(container):
    for ph in container.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            return ph
    return None


# --- placeholder surgery -----------------------------------------------------
# A table or a chart cannot hold text the way a body placeholder does, but it can
# still BE the placeholder: PowerPoint, when you insert a table into a content
# placeholder, keeps the <p:ph> marker on the resulting graphicFrame. We do the
# same — move the marker onto the new frame and drop the emptied placeholder — so
# the shape stays master-governed (it inherits the layout's position and moves
# with it) instead of becoming a free object floated onto the slide.
def _adopt_placeholder(ph, frame):
    """Give `frame` (table/chart graphicFrame) the placeholder identity of `ph`."""
    try:
        src = ph._element.nvSpPr.nvPr.find(qn("p:ph"))
        if src is None:
            return
        nv_pr = frame._element.nvGraphicFramePr.nvPr
        for old in nv_pr.findall(qn("p:ph")):
            nv_pr.remove(old)
        nv_pr.insert(0, src)          # <p:ph> is the first child of <p:nvPr>
        ph._element.getparent().remove(ph._element)
    except Exception:                 # never fail a build over the marker
        pass


def _drop_placeholder(ph):
    """Remove an unused placeholder so it does not sit empty on the slide."""
    if ph is None:
        return
    try:
        ph._element.getparent().remove(ph._element)
    except Exception:
        pass


def _prep_ph_tf(ph, anchor=MSO_ANCHOR.TOP, clear=True):
    """Ready a placeholder's text frame for our content while keeping it a REAL
    placeholder (inherits the layout/master; moves when the layout moves). Autofit
    is turned OFF so the template's shrink-to-fit can never override our floors."""
    tf = ph.text_frame
    if clear:
        tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    return tf


def setup_layouts(prs, theme, g):
    """Configure — ONCE, on the layouts — the geometry of every placeholder the
    default renderer writes into. Slides created from a layout INHERIT this, so
    titles AND body content are master-governed: reposition a placeholder on the
    layout in PowerPoint and every slide built on it moves together."""
    gutter = 0.5
    col_w = (g["contentW"] - gutter) / 2

    # Title and Content: bottom-anchored title region + body region on the grid.
    ct, cbodies = _phs_by_role(prs.slide_layouts[CONTENT_LAYOUT])
    _place(ct, g["marginX"], g["top"], g["contentW"], g["titleH"])
    ct.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM   # two-line titles grow upward
    if cbodies:
        _place(cbodies[0], g["marginX"], g["bodyTop"], g["contentW"], g["bodyH"])
        cbodies[0].text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Quote: its own layout ("Title Only"), so the centered quote block is layout
    # geometry like every other family — never a slide-level override.
    qt, _ = _phs_by_role(prs.slide_layouts[QUOTE_LAYOUT])
    if qt is not None:
        _place(qt, g["marginX"], 2.4, g["contentW"], 2.2)
        qt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title Slide: title block + subtitle below it.
    ts, tsb = _phs_by_role(prs.slide_layouts[TITLE_LAYOUT])
    _place(ts, g["marginX"], 2.45, g["contentW"], 1.3)
    if tsb:
        _place(tsb[0], g["marginX"], 3.78, g["contentW"], 0.8)

    # Section Header: an eyebrow/number body above, the section title below.
    sc, scb = _phs_by_role(prs.slide_layouts[SECTION_LAYOUT])
    _place(sc, g["marginX"], 2.86, g["contentW"], 1.4)
    if scb:
        _place(scb[0], g["marginX"], 2.18, g["contentW"], 0.6)
        scb[0].text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

    # Two Content: title region + two body columns on the grid.
    tc, tcb = _phs_by_role(prs.slide_layouts[TWO_CONTENT_LAYOUT])
    _place(tc, g["marginX"], g["top"], g["contentW"], g["titleH"])
    tc.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    for idx, ph in enumerate(tcb[:2]):
        _place(ph, g["marginX"] + idx * (col_w + gutter), g["bodyTop"], col_w, g["bodyH"])
        ph.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Picture with Caption: title region, image region, caption band beneath it.
    im = prs.slide_layouts[IMAGE_LAYOUT]
    it, itb = _phs_by_role(im)
    ipic = _picture_ph(im)
    _place(it, g["marginX"], g["top"], g["contentW"], g["titleH"])
    it.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    cap_top = g["footY"] - 1.0
    if ipic is not None:
        _place(ipic, g["marginX"], g["bodyTop"], g["contentW"], cap_top - g["bodyTop"] - 0.16)
    if itb:
        _place(itb[0], g["marginX"], cap_top, g["contentW"], 0.9)
        itb[0].text_frame.vertical_anchor = MSO_ANCHOR.TOP


def _fit_title_size(text, theme, g, width=None):
    """Largest size in [title_min, title_max] that keeps the title to two lines.
    Returns (size, overflow) — overflow True means it still needs >2 lines at the
    floor, so we keep it at the floor (never tiny) and warn to split/shorten."""
    sz, w = theme["size"], (width or g["contentW"])
    for size in range(int(sz["title_max"]), int(sz["title_min"]) - 1, -1):
        if _est_lines(text, size, w) <= 2:
            return size, False
    return int(sz["title_min"]), True


def _set_bg(slide, theme):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(theme["color"]["bg"])


def _hairline(slide, theme, g, y):
    if not theme.get("rule"):
        return
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(g["marginX"]), Inches(y),
                                 Inches(g["ruleLen"]), Inches(g["ruleH"]))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(theme["color"]["accent"])
    shp.line.fill.background()
    shp.shadow.inherit = False


def _title(slide, theme, g, text, idx):
    """Write the title into the inherited layout title placeholder, auto-fit size."""
    text = text or ""
    size, overflow = _fit_title_size(text, theme, g)
    if overflow:
        _warn("slide %d: title needs >2 lines even at %dpt — shorten or split: %r"
              % (idx, int(theme["size"]["title_min"]), text[:24] + ("…" if len(text) > 24 else "")))
    ph = slide.shapes.title
    tf = ph.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.12
    _no_bullet(p)
    run = p.add_run()
    run.text = text
    _set_run_font(run, name=theme["font"]["heading"], size=size, bold=True,
                  color=theme["color"]["ink"])
    _hairline(slide, theme, g, g["ruleY"]["content"])


def _fill_bullets(tf, theme, items):
    sz = theme["size"]
    first = True
    for it in _normalize_bullets(items):
        lvl = it["level"]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _bullet_para(p, "–" if lvl == 0 else "•", lvl, theme["font"]["body"])
        run = p.add_run()
        run.text = it["text"]
        _set_run_font(run, name=theme["font"]["body"],
                      size=sz["body"] if lvl == 0 else sz["body_sub"],
                      color=theme["color"]["ink" if lvl == 0 else "muted"])


def _bullets_height(items, theme, width, heading=None):
    """Estimate the rendered height (inches) of a bullet column at readable sizes."""
    sz = theme["size"]
    h = 0.0
    if heading:
        h += _est_lines(heading, sz["body"], width) * (sz["body"] * 1.32 / 72) + 8 / 72
    for it in _normalize_bullets(items):
        size = sz["body"] if it["level"] == 0 else sz["body_sub"]
        avail = width - (0.30 + 0.22 * it["level"])
        h += _est_lines(it["text"], size, avail) * (size * 1.34 / 72) + 9 / 72
    return h


def _check_body_overflow(height_in, g, idx):
    if height_in > g["bodyH"] + 0.05:
        _warn("slide %d: body needs ~%.1f in but only %.1f in fits at readable sizes"
              " — split into two slides (one message per slide)" % (idx, height_in, g["bodyH"]))


def _source(slide, theme, g, text):
    if not text:
        return
    tf = add_textbox(slide, g["marginX"], g["footY"] - 0.02, g["contentW"] - 0.9, 0.3)
    set_simple(tf, text, theme, font="body", size=theme["size"]["source"], color="muted")


def _page_number(slide, theme, g, n):
    if not theme.get("pageNumbers"):
        return
    tf = add_textbox(slide, g["pageW"] - g["marginX"] - 0.9, g["footY"], 0.9, 0.3)
    set_simple(tf, str(n), theme, font="body", size=theme["size"]["page_number"],
               color="muted", align=PP_ALIGN.RIGHT)


def _fill_col(tf, theme, col):
    """Write a column's heading (accent) + bullets into a placeholder text frame."""
    items = _normalize_bullets(col.get("bullets"))
    first = True
    if col.get("heading"):
        p = tf.paragraphs[0]
        _no_bullet(p)
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = col["heading"]
        _set_run_font(run, name=theme["font"]["heading"], size=theme["size"]["body"],
                      bold=True, color=theme["color"]["accent"])
        first = False
    for it in items:
        lvl = it["level"]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _bullet_para(p, "–" if lvl == 0 else "•", lvl, theme["font"]["body"])
        run = p.add_run()
        run.text = it["text"]
        _set_run_font(run, name=theme["font"]["body"],
                      size=theme["size"]["body"] if lvl == 0 else theme["size"]["body_sub"],
                      color=theme["color"]["ink" if lvl == 0 else "muted"])


def _trailing_para(tf, text, theme, font, size, color, space_before=6):
    """Add a non-bulleted trailing line (attribution, big-number caption)."""
    p = tf.add_paragraph()
    _no_bullet(p)
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    _set_run_font(run, name=theme["font"][font], size=size, color=theme["color"][color])
    return p


# ---------------------------------------------------------------------------
# COMPOSED PARTS — the closed vocabulary.
#
# The rule that keeps this from becoming decoration: a part is only ever drawn
# because the CONTENT has that shape (three parallel units -> cards; a real
# sequence -> steps; two meaningful axes -> matrix). Nothing here is available
# as an ornament, and every part carries a name (`part/<kind>`) so `audit_pptx.py`
# can tell a composed graphic from a textbox someone floated onto the slide.
#
# Geometry always comes from the BODY PLACEHOLDER's region: the layout still
# decides where content lives, the part just fills that region and the emptied
# placeholder is dropped (the same contract `image` slides use).
# ---------------------------------------------------------------------------
def _region(slide, theme, g, i, drop=True):
    """(x, y, w, h) of the body placeholder, in inches; the placeholder is then
    dropped so it cannot sit empty behind the part."""
    _, bodies = _phs_by_role(slide)
    if not bodies:
        _warn("slide %d: layout has no body placeholder — the part was placed on "
              "the content grid instead" % i)
        return g["marginX"], g["bodyTop"], g["contentW"], g["bodyH"]
    ph = bodies[0]
    r = (Emu(ph.left).inches, Emu(ph.top).inches,
         Emu(ph.width).inches, Emu(ph.height).inches)
    if drop:
        _drop_placeholder(ph)
    return r


def _part(slide, kind, x, y, w, h, theme, fill=None, line=None):
    """One part of the vocabulary: a rectangle with the deck's single radius, its
    single line weight, no shadow, and a name that records what it is."""
    radius = theme["shape"]["radius"]
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.name = "part/%s" % kind
    if radius > 0:
        try:    # the adjustment is the radius as a fraction of the shorter side
            sp.adjustments[0] = max(0.0, min(0.5, radius / max(0.01, min(w, h))))
        except Exception:
            pass
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(theme["shape"]["line"])
    else:
        sp.line.fill.background()
    try:
        sp.shadow.inherit = False        # no drop shadow, ever
    except Exception:
        pass
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    pad = Inches(_u(theme, 2))
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = pad
    return sp


def _part_text(sp, theme, label, text, index=None, accent=False):
    """Label (+ optional index) then explanation, inside a part. Every paragraph
    is aligned LEFT explicitly: text in an autoshape defaults to centered, and a
    deck whose cards are centered while its bullets are not reads as careless."""
    tf = sp.text_frame
    tf.clear()
    c = theme["color"]
    first = True
    if index is not None:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _no_bullet(p)
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = str(index)
        _set_run_font(run, name=theme["font"]["heading"], size=theme["size"]["part_index"],
                      bold=True, color=c["accent"])
        first = False
    if label:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        _no_bullet(p)
        p.space_after = Pt(5)
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = str(label)
        _set_run_font(run, name=theme["font"]["heading"], size=theme["size"]["part_label"],
                      bold=True, color=c["accent"] if accent else c["ink"])
        first = False
    if text:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        _no_bullet(p)
        p.line_spacing = 1.25
        run = p.add_run()
        run.text = str(text)
        _set_run_font(run, name=theme["font"]["body"], size=theme["size"]["part_text"],
                      color=c["muted"])


def _part_items(s, key):
    """[{label, text}] from a list of strings or dicts — or from a single dict,
    which is how a one-item field like `lead:` is naturally written."""
    v = s.get(key)
    if isinstance(v, dict):
        return [{"label": v.get("label", ""), "text": v.get("text", "")}]
    out = []
    for it in v or []:
        if isinstance(it, str):
            out.append({"label": it, "text": ""})
        else:
            out.append({"label": it.get("label", ""), "text": it.get("text", "")})
    return out


def _part_height(items, theme, w, min_h, max_h):
    """Tall enough for the wordiest item at readable sizes, within the region."""
    pad = _u(theme, 4)
    inner = max(0.5, w - pad)
    need = 0.0
    for it in items:
        lines_l = _est_lines(it["label"], theme["size"]["part_label"], inner)
        lines_t = _est_lines(it["text"], theme["size"]["part_text"], inner)
        need = max(need, lines_l * theme["size"]["part_label"] * 1.25 / 72.0
                   + lines_t * theme["size"]["part_text"] * 1.35 / 72.0)
    return max(min_h, min(max_h, need + pad))


def _render_cards(slide, theme, g, s, i):
    """2-4 EQUIVALENT units, side by side. Equivalent is the condition: if the
    items are not the same kind of thing, they are a list, not cards."""
    items = _part_items(s, "cards")
    if not 2 <= len(items) <= 4:
        _warn("slide %d: cards takes 2-4 equivalent items (got %d) — use bullets, or split"
              % (i, len(items)))
        if not items:
            return
    x, y, w, h = _region(slide, theme, g, i)
    n = len(items)
    gap = _u(theme, 2)
    cw = (w - gap * (n - 1)) / n
    ch = _part_height(items, theme, cw, min(1.3, h), min(h, 3.0))
    y += max(0.0, (h - ch) / 2)          # the row sits centered in its region
    emph = s.get("emphasis")
    c = theme["color"]
    for k, it in enumerate(items):
        hot = (emph is not None and int(emph) == k + 1)
        sp = _part(slide, "card", x + k * (cw + gap), y, cw, ch, theme,
                   fill=c["surface_hi"] if hot else c["surface"],
                   line=c["accent"] if hot else c["border"])
        _part_text(sp, theme, it["label"], it["text"], accent=hot)


def _render_steps(slide, theme, g, s, i):
    """A real sequence, 3-5 stages. This is the ONE archetype allowed to draw
    arrows, because here the arrow encodes the order that the content has."""
    items = _part_items(s, "steps")
    if not 3 <= len(items) <= 5:
        _warn("slide %d: steps takes 3-5 stages (got %d) — beyond that it is a process "
              "diagram, not a slide" % (i, len(items)))
        if not items:
            return
    x, y, w, h = _region(slide, theme, g, i)
    n = len(items)
    arrow_w, arrow_gap = _u(theme, 3), _u(theme, 1)
    span = arrow_w + arrow_gap * 2
    cw = (w - span * (n - 1)) / n
    ch = _part_height(items, theme, cw, min(1.3, h), min(h, 2.6))
    y += max(0.0, (h - ch) / 2)
    c = theme["color"]
    for k, it in enumerate(items):
        left = x + k * (cw + span)
        sp = _part(slide, "step", left, y, cw, ch, theme,
                   fill=c["surface"], line=c["border"])
        _part_text(sp, theme, it["label"], it["text"], index="%02d" % (k + 1))
        if k < n - 1:
            # A small triangle, not Office's block arrow: the mark has to say
            # "then" without becoming the loudest thing on the slide.
            tip = _u(theme, 1.6)
            ar = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                Inches(left + cw + arrow_gap + (arrow_w - tip) / 2),
                Inches(y + ch / 2 - tip / 2), Inches(tip), Inches(tip))
            ar.rotation = 90
            ar.name = "part/arrow"
            ar.fill.solid()
            ar.fill.fore_color.rgb = RGBColor.from_string(c["muted"])
            ar.line.fill.background()
            try:
                ar.shadow.inherit = False
            except Exception:
                pass


def _render_lead(slide, theme, g, s, i):
    """One unit that matters more, beside 2-3 that support it. The only archetype
    where the parts are deliberately UNEQUAL: `cards` says "compare these", this
    says "this one, and the others are context". Area is the argument, so the
    lead gets ~46% of the width and the full height."""
    lead = _part_items(s, "lead")
    rest = _part_items(s, "rest")
    if not lead:
        _warn("slide %d: lead slide has no `lead:` item — nothing to feature" % i)
        return
    if not 2 <= len(rest) <= 3:
        _warn("slide %d: lead takes 2-3 supporting items (got %d)" % (i, len(rest)))
        if not rest:
            return
    x, y, w, h = _region(slide, theme, g, i)
    c = theme["color"]
    gap = _u(theme, 2)
    lw = (w - gap) * 0.46
    rw = (w - gap) * 0.54
    sp = _part(slide, "lead", x, y, lw, h, theme,
               fill=c["surface_hi"], line=c["accent"])
    _part_text(sp, theme, lead[0]["label"], lead[0]["text"], accent=True)
    n = len(rest)
    rh = (h - gap * (n - 1)) / n
    for k, it in enumerate(rest):
        sp = _part(slide, "rest", x + lw + gap, y + k * (rh + gap), rw, rh, theme,
                   fill=c["surface"], line=c["border"])
        _part_text(sp, theme, it["label"], it["text"])


def _render_matrix(slide, theme, g, s, i):
    """Two axes that BOTH carry meaning, and four quadrants that all say something.
    If one quadrant is empty filler, the content is a list wearing a matrix."""
    items = _part_items(s, "quadrants")
    if len(items) != 4:
        _warn("slide %d: matrix needs exactly 4 quadrants (top-left, top-right, "
              "bottom-left, bottom-right); got %d" % (i, len(items)))
        if not items:
            return
        items = (items + [{"label": "", "text": ""}] * 4)[:4]
    x, y, w, h = _region(slide, theme, g, i)
    c = theme["color"]
    xa = s.get("x_axis") or []
    ya = s.get("y_axis") or []
    gut = _u(theme, 14) if ya else 0.0     # left gutter: a real word has to fit
    foot = _u(theme, 4) if xa else 0.0     # bottom strip for the x-axis labels
    x0, y0 = x + gut, y
    w0, h0 = w - gut, h - foot
    gap = _u(theme, 1.5)
    cw, ch = (w0 - gap) / 2, (h0 - gap) / 2
    emph = s.get("emphasis")
    for k, it in enumerate(items):
        col, row = k % 2, k // 2
        hot = (emph is not None and int(emph) == k + 1)
        sp = _part(slide, "quadrant", x0 + col * (cw + gap), y0 + row * (ch + gap),
                   cw, ch, theme,
                   fill=c["surface_hi"] if hot else c["surface"],
                   line=c["accent"] if hot else c["border"])
        _part_text(sp, theme, it["label"], it["text"], accent=hot)
    sz = theme["size"]["axis"]
    if len(ya) >= 2:      # vertical axis: high at the top, low at the bottom
        for text, ty in ((ya[1], y0), (ya[0], y0 + h0 - _u(theme, 3))):
            tf = add_textbox(slide, x, ty, gut - _u(theme, 2), _u(theme, 3),
                             name="part/axis")
            set_simple(tf, text, theme, font="body", size=sz, color="muted",
                       align=PP_ALIGN.RIGHT)
    if len(xa) >= 2:      # horizontal axis, in the strip under the quadrants
        for text, tx, al in ((xa[0], x0, PP_ALIGN.LEFT),
                             (xa[1], x0 + w0 / 2, PP_ALIGN.RIGHT)):
            tf = add_textbox(slide, tx, y0 + h0 + _u(theme, 0.5), w0 / 2, _u(theme, 3),
                             name="part/axis")
            set_simple(tf, text, theme, font="body", size=sz, color="muted", align=al)


def _render_split(slide, theme, g, s, i):
    """The asymmetric composition: figure and its reading, 62/38. Symmetry says
    'these are equal'; most figure-plus-explanation slides are not."""
    img = _asset(s.get("image"))
    _, bodies = _phs_by_role(slide)
    if not bodies:
        _warn("slide %d: layout has no body placeholder — split not drawn" % i)
        return
    ph = bodies[0]
    x, y = Emu(ph.left).inches, Emu(ph.top).inches
    w, h = Emu(ph.width).inches, Emu(ph.height).inches
    # 38/62 by default. A WIDE figure wants more of the slide (0.3), a tall one
    # less (0.45) — `ratio` is the text column's share, clamped to a range that
    # keeps both halves usable.
    ratio = s.get("ratio", 0.38)
    try:
        ratio = max(0.25, min(0.50, float(ratio)))
    except (TypeError, ValueError):
        ratio = 0.38
    gap = _u(theme, 4)
    text_w = (w - gap) * ratio
    fig_w = (w - gap) * (1.0 - ratio)
    flip = bool(s.get("flip"))
    text_x = x if not flip else x + fig_w + gap
    fig_x = x + text_w + gap if not flip else x
    _place(ph, text_x, y, text_w, h)
    ph.name = "part/split-text"      # records that this geometry IS the composition
    _fill_col(_prep_ph_tf(ph), theme, {"heading": s.get("heading"),
                                       "bullets": s.get("bullets")})
    if img and os.path.exists(img):
        pic = slide.shapes.add_picture(img, Inches(fig_x), Inches(y), height=Inches(h))
        if pic.width > Inches(fig_w):
            pic.height = int(pic.height * Inches(fig_w) / pic.width)
            pic.width = Inches(fig_w)
        pic.left = Inches(fig_x + (fig_w - Emu(pic.width).inches) / 2)
        pic.top = Inches(y + (h - Emu(pic.height).inches) / 2)
    else:
        _warn("slide %d: image not found (%r) — the figure half is empty"
              % (i, s.get("image")))
        sp = _part(slide, "figure-missing", fig_x, y, fig_w, h, theme,
                   fill=theme["color"]["surface"], line=theme["color"]["border"])
        _part_text(sp, theme, "", "[ image: %s ]" % (s.get("image") or "missing"))


def _render_statement(slide, theme, g, s, i):
    """One sentence, alone. The deck's punctuation: a turn, a verdict, a stake in
    the ground. Never a slide that merely has little on it."""
    title_ph, _ = _phs_by_role(slide)
    if title_ph is None:
        return
    # Anchored TOP so the sentence sits directly under the rule: the rule is the
    # mark that something is being declared, and a gap between them breaks that.
    tf = _prep_ph_tf(title_ph, anchor=MSO_ANCHOR.TOP)
    set_simple(tf, s.get("text") or s.get("title") or "", theme, font="heading",
               size=theme["size"]["statement"], bold=True, color="ink",
               align=PP_ALIGN.LEFT, line_spacing=1.25)
    if s.get("sub"):
        p = _trailing_para(tf, s["sub"], theme, "body", theme["size"]["statement_sub"],
                           "muted", space_before=12)
        if p is not None:
            p.alignment = PP_ALIGN.LEFT


# ---------------------------------------------------------------------------
# Tables and charts. Both exist so a deck that needs one does NOT have to fall
# back to hand-written python-pptx (which is how decks end up as free textboxes
# on blank slides). Both are inserted at the BODY PLACEHOLDER's geometry and
# adopt its placeholder marker, so they stay master-governed.
# ---------------------------------------------------------------------------
PLAIN_TABLE_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"   # "No Style, No Grid"
_LN_ORDER = ("a:lnL", "a:lnR", "a:lnT", "a:lnB", "a:lnTlToBr", "a:lnBlToTr")


def _cell_border(cell, edge, color, pt):
    """Draw one hairline edge on a table cell (python-pptx has no border API)."""
    tag = "a:" + edge
    tcPr = cell._tc.get_or_add_tcPr()
    for el in tcPr.findall(qn(tag)):
        tcPr.remove(el)
    ln = tcPr.makeelement(qn(tag), {"w": str(int(pt * 12700)), "cap": "flat",
                                    "cmpd": "sng", "algn": "ctr"})
    fill = ln.makeelement(qn("a:solidFill"), {})
    fill.append(ln.makeelement(qn("a:srgbClr"), {"val": color}))
    ln.append(fill)
    want = _LN_ORDER.index(tag)
    for child in list(tcPr):
        name = "a:" + child.tag.split("}")[-1]
        if name not in _LN_ORDER or _LN_ORDER.index(name) > want:
            child.addprevious(ln)
            return
    tcPr.append(ln)


def _table_rows(s):
    """Normalize the spec into (header list or None, [row, ...]) of plain strings."""
    header = s.get("columns") or s.get("header")
    header = [str(c) for c in header] if header else None
    rows = []
    for r in s.get("rows") or []:
        if isinstance(r, dict) and header:
            rows.append([str(r.get(c, "")) for c in header])
        elif isinstance(r, (list, tuple)):
            rows.append([str(c) for c in r])
        else:
            rows.append([str(r)])
    return header, rows


def _render_table(slide, theme, g, s, i):
    header, rows = _table_rows(s)
    if not header and not rows:
        _warn("slide %d: table has no columns and no rows — nothing to draw" % i)
        return
    ncols = max([len(header or [])] + [len(r) for r in rows] or [1]) or 1
    nrows = len(rows) + (1 if header else 0)
    if ncols > 6:
        _warn("slide %d: table has %d columns — 6 is the readable maximum, split it"
              % (i, ncols))
    _, bodies = _phs_by_role(slide)
    if not bodies:
        _warn("slide %d: layout has no body placeholder — table not drawn" % i)
        return
    ph = bodies[0]
    left, top, width = ph.left, ph.top, ph.width
    height = min(ph.height, Inches(max(0.36 * nrows, 0.4)))
    if Emu(int(height)).inches > g["bodyH"] + 0.05:
        _warn("slide %d: table needs ~%.1f in but only %.1f in fits at readable sizes"
              " — split the table across slides" % (i, Emu(int(height)).inches, g["bodyH"]))
    frame = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    _adopt_placeholder(ph, frame)
    tbl = frame.table

    widths = s.get("widths")
    if widths and len(widths) == ncols and sum(widths) > 0:
        total = float(sum(widths))
        for c, w in enumerate(widths):
            tbl.columns[c].width = Emu(int(width * (w / total)))

    tbl.first_row = bool(header)
    tbl.horz_banding = False
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        for el in tblPr.findall(qn("a:tableStyleId")):
            tblPr.remove(el)
        sid = tblPr.makeelement(qn("a:tableStyleId"), {})
        sid.text = PLAIN_TABLE_STYLE
        tblPr.append(sid)

    sz = theme["size"]
    body_rows = ([header] if header else []) + rows
    for r, data in enumerate(body_rows):
        is_head = bool(header) and r == 0
        tbl.rows[r].height = Emu(int(Inches(0.36)))
        for c in range(ncols):
            cell = tbl.cell(r, c)
            cell.fill.background()
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _no_bullet(p)
            run = p.add_run()
            run.text = data[c] if c < len(data) else ""
            _set_run_font(run, name=theme["font"]["heading" if is_head else "body"],
                          size=sz["table_header"] if is_head else sz["table"],
                          bold=is_head,
                          color=theme["color"]["ink" if is_head else "ink"])
            # header: accent rule under it; body rows: hairline separator
            _cell_border(cell, "lnB", theme["color"]["accent" if is_head else "muted"],
                         1.25 if is_head else 0.5)


_CHART_TYPES = {
    "column": "COLUMN_CLUSTERED", "bar": "BAR_CLUSTERED",
    "line": "LINE_MARKERS", "area": "AREA", "pie": "PIE", "doughnut": "DOUGHNUT",
}


def _chart_no_border(chart):
    """Drop the chart-area frame (the default box around the plot) and its fill."""
    try:
        cs = chart._chartSpace
        for el in cs.findall(qn("c:spPr")):
            cs.remove(el)
        spPr = cs.makeelement(qn("c:spPr"), {})
        spPr.append(spPr.makeelement(qn("a:noFill"), {}))
        ln = spPr.makeelement(qn("a:ln"), {})
        ln.append(ln.makeelement(qn("a:noFill"), {}))
        spPr.append(ln)
        txPr = cs.find(qn("c:txPr"))
        (txPr.addprevious(spPr) if txPr is not None else cs.append(spPr))
    except Exception:
        pass


def _series_colors(theme, style="focus"):
    """focus: the series that carries the message in accent, the rest in greys
    that recede (light to dark, so they never compete with each other).
    tonal: one hue at four lightnesses — right when the series are the SAME kind
    of thing and the comparison is between them, not against one of them."""
    pal = theme.get("series")
    if pal:
        return [str(c).lstrip("#") for c in pal]
    accent = theme["color"]["accent"]
    if str(style).lower() == "tonal":
        return [accent if t is None else _tone(accent, t) for t in SERIES_TONES]
    return [accent, "B6BDC8", "8A93A1", "5B6472", "343B45"]


def _render_chart(slide, theme, g, s, i):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    kind = str(s.get("chart", "column")).lower()
    if kind not in _CHART_TYPES:
        _warn("slide %d: unknown chart type %r — using 'column' "
              "(column/bar/line/area/pie/doughnut)" % (i, kind))
        kind = "column"
    cats = [str(c) for c in (s.get("categories") or [])]
    series = s.get("series") or []
    if isinstance(series, dict):          # {name: [values]} shorthand
        series = [{"name": k, "values": v} for k, v in series.items()]
    series = [se for se in series if se.get("values")]
    if not cats or not series:
        _warn("slide %d: chart needs `categories` and at least one `series` — skipped" % i)
        return
    _, bodies = _phs_by_role(slide)
    if not bodies:
        _warn("slide %d: layout has no body placeholder — chart not drawn" % i)
        return
    ph = bodies[0]

    data = CategoryChartData()
    data.categories = cats
    for se in series:
        vals = list(se.get("values") or [])
        vals += [None] * (len(cats) - len(vals))
        data.add_series(str(se.get("name", "")), vals[:len(cats)])

    frame = slide.shapes.add_chart(getattr(XL_CHART_TYPE, _CHART_TYPES[kind]),
                                   ph.left, ph.top, ph.width, ph.height, data)
    _adopt_placeholder(ph, frame)
    chart = frame.chart
    _chart_no_border(chart)
    chart.has_title = False     # the SLIDE title states the message; a chart title repeats it
    sz = theme["size"]
    chart.font.size = Pt(sz["chart_label"])
    chart.font.color.rgb = RGBColor.from_string(theme["color"]["muted"])
    chart.font.name = theme["font"]["body"]

    legend = s.get("legend")
    show_legend = (len(series) > 1) if legend is None else bool(legend)
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    pal = _series_colors(theme, s.get("series_style", theme.get("series_style", "focus")))
    plot = chart.plots[0]
    plot.gap_width = 60
    if kind in ("pie", "doughnut"):
        plot.vary_by_categories = True
        pts = list(plot.series[0].points)
        for idx, pt in enumerate(pts):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = RGBColor.from_string(pal[idx % len(pal)])
        chart.has_legend = True if legend is None else bool(legend)
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
    else:
        plot.vary_by_categories = False
        for idx, ser in enumerate(plot.series):
            color = RGBColor.from_string(pal[idx % len(pal)])
            if kind in ("line",):
                ser.format.line.color.rgb = color
                ser.format.line.width = Pt(2.25)
            else:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = color
                ser.format.line.fill.background()
        try:                                   # quiet axes: no gridlines, hairline base
            va = chart.value_axis
            va.has_major_gridlines = bool(s.get("gridlines", False))
            if va.has_major_gridlines:
                va.major_gridlines.format.line.color.rgb = \
                    RGBColor.from_string(theme["color"]["muted"])
                va.major_gridlines.format.line.width = Pt(0.5)
            va.format.line.color.rgb = RGBColor.from_string(theme["color"]["muted"])
            ca = chart.category_axis
            ca.format.line.color.rgb = RGBColor.from_string(theme["color"]["muted"])
            ca.has_major_gridlines = False
        except Exception:
            pass
    if s.get("data_labels"):
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(sz["chart_label"])
        plot.data_labels.font.color.rgb = RGBColor.from_string(theme["color"]["ink"])


def render_default(prs, theme, g, slides):
    """Every slide is built on a STANDARD layout and its content is written into
    that layout's real placeholders — no free textboxes floated onto blank pages.
    (Peripheral footnotes — source line, page number — remain small annotations.)"""
    setup_layouts(prs, theme, g)
    for i, s in enumerate(slides, start=1):
        t = s.get("type", "bullets")

        if t == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[TITLE_LAYOUT])
            _set_bg(slide, theme)
            _hairline(slide, theme, g, g["ruleY"]["title"])
            title_ph, bodies = _phs_by_role(slide)
            _prep_ph_tf(title_ph)
            set_simple(title_ph.text_frame, s.get("title", ""), theme, font="heading",
                       size=theme["size"]["title_slide"], bold=True, color="ink", line_spacing=1.1)
            if s.get("subtitle") and bodies:
                _prep_ph_tf(bodies[0])
                set_simple(bodies[0].text_frame, s["subtitle"], theme,
                           font="body", size=theme["size"]["subtitle"], color="muted")
            continue

        if t == "section":
            th = _inverted(theme) if s.get("invert") else theme
            slide = prs.slides.add_slide(prs.slide_layouts[SECTION_LAYOUT])
            _set_bg(slide, th)
            _hairline(slide, th, g, g["ruleY"]["section"])
            title_ph, bodies = _phs_by_role(slide)
            if s.get("number") is not None and bodies:
                _prep_ph_tf(bodies[0], anchor=MSO_ANCHOR.BOTTOM)
                set_simple(bodies[0].text_frame, str(s["number"]), th, font="heading",
                           size=th["size"]["section_number"], bold=True, color="accent")
            _prep_ph_tf(title_ph)
            set_simple(title_ph.text_frame, s.get("title", ""), th, font="heading",
                       size=th["size"]["section"], bold=True, color="ink", line_spacing=1.1)
            continue

        if t == "quote":
            # "Title Only" — the quote lives in that layout's TITLE placeholder, whose
            # geometry was set once in setup_layouts (no slide-level override).
            slide = prs.slides.add_slide(prs.slide_layouts[QUOTE_LAYOUT])
            _set_bg(slide, theme)
            _hairline(slide, theme, g, g["ruleY"]["quote"])
            quote_ph, _ = _phs_by_role(slide)
            if quote_ph is not None:
                tf = _prep_ph_tf(quote_ph, anchor=MSO_ANCHOR.MIDDLE)
                set_simple(tf, "“" + s.get("quote", "") + "”", theme, font="heading",
                           size=theme["size"]["quote"], color="ink", line_spacing=1.25)
                if s.get("attribution"):
                    _trailing_para(tf, "— " + s["attribution"], theme, "body",
                                   theme["size"]["quote_attr"], "muted", space_before=10)
            continue

        if t == "two_col":
            slide = prs.slides.add_slide(prs.slide_layouts[TWO_CONTENT_LAYOUT])
            _set_bg(slide, theme)
            _title(slide, theme, g, s.get("title", ""), i)
            _, bodies = _phs_by_role(slide)
            gutter = 0.5
            col_w = (g["contentW"] - gutter) / 2
            for idx, key in enumerate(("left", "right")):
                if idx >= len(bodies):
                    break
                col = s.get(key) or {}
                _check_body_overflow(
                    _bullets_height(col.get("bullets"), theme, col_w, col.get("heading")), g, i)
                _fill_col(_prep_ph_tf(bodies[idx]), theme, col)
            _source(slide, theme, g, s.get("source"))
            _page_number(slide, theme, g, i)
            continue

        if t == "image":
            slide = prs.slides.add_slide(prs.slide_layouts[IMAGE_LAYOUT])
            _set_bg(slide, theme)
            _title(slide, theme, g, s.get("title", ""), i)
            _render_image(slide, theme, g, s, i)
            continue

        if t == "statement":
            # "Title Only" — one sentence in that layout's title placeholder.
            th = _inverted(theme) if s.get("invert") else theme
            slide = prs.slides.add_slide(prs.slide_layouts[QUOTE_LAYOUT])
            _set_bg(slide, th)
            _hairline(slide, th, g, g["ruleY"]["quote"])
            _render_statement(slide, th, g, s, i)
            _page_number(slide, th, g, i)
            continue

        if t in ("cards", "steps", "matrix", "split", "lead"):
            slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_LAYOUT])
            _set_bg(slide, theme)
            _title(slide, theme, g, s.get("title", ""), i)
            {"cards": _render_cards, "steps": _render_steps, "lead": _render_lead,
             "matrix": _render_matrix, "split": _render_split}[t](slide, theme, g, s, i)
            _source(slide, theme, g, s.get("source"))
            _page_number(slide, theme, g, i)
            continue

        if t == "big_number":
            slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_LAYOUT])
            _set_bg(slide, theme)
            _title(slide, theme, g, s.get("title", ""), i)
            _, bodies = _phs_by_role(slide)
            if bodies:
                tf = _prep_ph_tf(bodies[0], anchor=MSO_ANCHOR.MIDDLE)
                set_simple(tf, str(s.get("number", "")), theme, font="number",
                           size=theme["size"]["big_number"], bold=True, color="accent")
                if s.get("caption"):
                    _trailing_para(tf, s["caption"], theme, "body",
                                   theme["size"]["big_caption"], "muted")
            _source(slide, theme, g, s.get("source"))
            _page_number(slide, theme, g, i)
            continue

        if t in ("table", "chart"):
            slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_LAYOUT])
            _set_bg(slide, theme)
            _title(slide, theme, g, s.get("title", ""), i)
            (_render_table if t == "table" else _render_chart)(slide, theme, g, s, i)
            _source(slide, theme, g, s.get("source"))
            _page_number(slide, theme, g, i)
            continue

        if t == "blank":
            slide = _blank_with_bg(prs, theme)
            _page_number(slide, theme, g, i)
            continue

        # bullets (default)
        slide = prs.slides.add_slide(prs.slide_layouts[CONTENT_LAYOUT])
        _set_bg(slide, theme)
        _title(slide, theme, g, s.get("title", ""), i)
        _, bodies = _phs_by_role(slide)
        _check_body_overflow(_bullets_height(s.get("bullets"), theme, g["contentW"]), g, i)
        if bodies:
            _fill_bullets(_prep_ph_tf(bodies[0]), theme, s.get("bullets"))
        _source(slide, theme, g, s.get("source"))
        _page_number(slide, theme, g, i)


def _blank_with_bg(prs, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    _set_bg(slide, theme)
    return slide


def _render_image(slide, theme, g, s, i):
    """Fit the picture inside the layout's PICTURE placeholder region (uncropped,
    aspect preserved) and write the caption into its BODY caption placeholder — the
    figure sits in the master's designated region, not free-floated on a blank slide."""
    pic_ph = _picture_ph(slide)
    _, bodies = _phs_by_role(slide)
    cap_ph = bodies[0] if bodies else None
    label, note = s.get("caption"), (s.get("note") or s.get("description"))
    img = _asset(s.get("image"))

    # The picture placeholder defines the (master-governed) image region. Read it,
    # then drop the empty placeholder so it never lingers behind the fitted picture.
    if pic_ph is not None:
        rx, ry = Emu(pic_ph.left).inches, Emu(pic_ph.top).inches
        rw, rh = Emu(pic_ph.width).inches, Emu(pic_ph.height).inches
    else:
        rx, ry, rw, rh = g["marginX"], g["bodyTop"], g["contentW"], g["bodyH"] * 0.7

    if img and os.path.exists(img):
        _drop_placeholder(pic_ph)   # only once the real picture takes its region
        pic = slide.shapes.add_picture(img, Inches(rx), Inches(ry), height=Inches(rh))
        if pic.width > Inches(rw):                  # too wide: refit by width, keep aspect
            pic.height = int(pic.height * Inches(rw) / pic.width)
            pic.width = Inches(rw)
        pic.left = Inches(rx + (rw - Emu(pic.width).inches) / 2)
        pic.top = Inches(ry + (rh - Emu(pic.height).inches) / 2)
    else:
        _warn("slide %d: image not found (%r) — caption kept, image region left empty"
              % (i, img))
        marker = "[ image: %s ]" % (img or "missing")
        try:    # keep the marker INSIDE the picture placeholder, not floating over it
            tf = _prep_ph_tf(pic_ph, anchor=MSO_ANCHOR.MIDDLE)
        except Exception:
            tf = add_textbox(slide, rx, ry, rw, rh, anchor=MSO_ANCHOR.MIDDLE)
        set_simple(tf, marker, theme, font="body", size=16, color="muted",
                   align=PP_ALIGN.CENTER)

    if cap_ph is not None and (label or note):
        tf = _prep_ph_tf(cap_ph)
        first = True
        if label:
            p = tf.paragraphs[0]
            _no_bullet(p)
            run = p.add_run()
            run.text = label
            _set_run_font(run, name=theme["font"]["heading"], size=theme["size"]["caption"],
                          bold=True, color=theme["color"]["ink"])
            first = False
        if note:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            _no_bullet(p)
            if not first:
                p.space_before = Pt(4)
            p.line_spacing = 1.15
            run = p.add_run()
            run.text = note
            _set_run_font(run, name=theme["font"]["body"], size=theme["size"]["caption_note"],
                          color=theme["color"]["muted"])
    _source(slide, theme, g, s.get("source"))
    _page_number(slide, theme, g, i)


# ---------------------------------------------------------------------------
# TEMPLATE-FILL MODE — open a real template, fill its placeholders.
# ---------------------------------------------------------------------------
def _ph_type_name(ph):
    try:
        return str(ph.placeholder_format.type).split(".")[-1].split(" ")[0]
    except Exception:
        return "?"


def _layout_index_by_role(prs):
    """Heuristic: pick a sensible layout index per slide type from layout names
    and placeholder types. Used when the spec/map does not pin layouts."""
    layouts = list(prs.slide_layouts)

    def find(pred, default=0):
        for i, lo in enumerate(layouts):
            if pred(lo):
                return i
        return default

    def has_types(lo, types):
        present = {ph.placeholder_format.type for ph in lo.placeholders}
        return any(t in present for t in types)

    def name_match(lo, *words):
        nm = (lo.name or "").lower()
        return any(w in nm for w in words)

    content = find(lambda lo: name_match(lo, "title and content", "content")
                   or has_types(lo, (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)), 1)
    # image: a layout with a real PICTURE placeholder wins over any "...caption" name.
    img = find(lambda lo: has_types(lo, (PP_PLACEHOLDER.PICTURE,)), -1)
    if img < 0:
        img = find(lambda lo: name_match(lo, "picture", "image", "caption"), content)
    return {
        "title": find(lambda lo: name_match(lo, "title slide")
                      or has_types(lo, (PP_PLACEHOLDER.SUBTITLE,)), 0),
        "section": find(lambda lo: name_match(lo, "section"), content),
        "two_col": find(lambda lo: name_match(lo, "two content", "comparison",
                                              "two-col", "2 content"), content),
        "bullets": content, "big_number": content, "quote": content,
        "table": content, "chart": content,
        # Composed archetypes are drawn into the body placeholder's region, so
        # they want the same layout a bullets slide would use.
        "cards": content, "steps": content, "matrix": content, "split": content,
        "lead": content,
        # A statement is one sentence alone: a title-only layout if the template
        # has one, else the section divider, else content.
        "statement": find(lambda lo: name_match(lo, "title only", "statement"),
                          find(lambda lo: name_match(lo, "section"), content)),
        "image": img,
        "blank": find(lambda lo: name_match(lo, "blank"), content),
    }


def _placeholders_by_role(slide):
    """Group a slide's placeholders into title / body[] / subtitle / picture by type."""
    roles = {"title": None, "subtitle": None, "body": [], "picture": None}
    for ph in slide.placeholders:
        t = ph.placeholder_format.type
        if t in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            roles["title"] = ph
        elif t == PP_PLACEHOLDER.SUBTITLE:
            roles["subtitle"] = ph
        elif t in (PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.OBJECT) and roles["picture"] is None \
                and t == PP_PLACEHOLDER.PICTURE:
            roles["picture"] = ph
        elif t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE):
            roles["body"].append(ph)
    # body keeps left-to-right reading order
    roles["body"].sort(key=lambda p: (p.top or 0, p.left or 0))
    return roles


def _ph_by_idx(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_ph_text(ph, text):
    """Write plain text into a placeholder, preserving its template formatting."""
    if ph is None or text is None:
        return
    ph.text_frame.text = str(text)


def _set_ph_bullets(ph, items):
    """Write multi-level bullets into a placeholder; the template styles them."""
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    norm = _normalize_bullets(items)
    if not norm:
        return
    for i, it in enumerate(norm):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it["text"]
        p.level = min(it["level"], 8)


def _tpl_need(ph, idx, layout_name, role, content):
    """Template mode must never drop content silently: warn when the chosen layout
    has no placeholder for a role the spec actually filled."""
    if ph is None and content not in (None, "", [], {}):
        _warn("slide %d: layout %r has no %s placeholder — that content was NOT written."
              " Pin a layout with `layout:` or map the role in --map." % (idx, layout_name, role))
    return ph


def _render_table_template(slide, s, i, layout_name):
    """Table in template mode: the template's own table style and theme fonts."""
    header, rows = _table_rows(s)
    if not header and not rows:
        return
    ncols = max([len(header or [])] + [len(r) for r in rows] or [1]) or 1
    nrows = len(rows) + (1 if header else 0)
    roles = _placeholders_by_role(slide)
    ph = roles["body"][0] if roles["body"] else None
    if _tpl_need(ph, i, layout_name, "body (for the table)", rows or header) is None:
        return
    frame = slide.shapes.add_table(nrows, ncols, ph.left, ph.top, ph.width,
                                   min(ph.height, Emu(int(Inches(max(0.36 * nrows, 0.4))))))
    _adopt_placeholder(ph, frame)
    tbl = frame.table
    tbl.first_row = bool(header)
    for r, data in enumerate(([header] if header else []) + rows):
        for c in range(ncols):
            tbl.cell(r, c).text_frame.text = data[c] if c < len(data) else ""


def _render_chart_template(slide, s, i, layout_name):
    """Chart in template mode: colors come from the template's theme."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    kind = str(s.get("chart", "column")).lower()
    kind = kind if kind in _CHART_TYPES else "column"
    cats = [str(c) for c in (s.get("categories") or [])]
    series = s.get("series") or []
    if isinstance(series, dict):
        series = [{"name": k, "values": v} for k, v in series.items()]
    series = [se for se in series if se.get("values")]
    if not cats or not series:
        _warn("slide %d: chart needs `categories` and at least one `series` — skipped" % i)
        return
    roles = _placeholders_by_role(slide)
    ph = roles["body"][0] if roles["body"] else None
    if _tpl_need(ph, i, layout_name, "body (for the chart)", series) is None:
        return
    data = CategoryChartData()
    data.categories = cats
    for se in series:
        vals = list(se.get("values") or [])
        vals += [None] * (len(cats) - len(vals))
        data.add_series(str(se.get("name", "")), vals[:len(cats)])
    frame = slide.shapes.add_chart(getattr(XL_CHART_TYPE, _CHART_TYPES[kind]),
                                   ph.left, ph.top, ph.width, ph.height, data)
    _adopt_placeholder(ph, frame)
    chart = frame.chart
    chart.has_title = False
    show = (len(series) > 1) if s.get("legend") is None else bool(s.get("legend"))
    chart.has_legend = show
    if show:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False


_THEME_REL = ("http://schemas.openxmlformats.org/officeDocument/"
              "2006/relationships/theme")


def _theme_of_template(prs, theme_path=DEFAULT_THEME):
    """Composed parts need colors, sizes and shape tokens, and in template-fill
    mode there is no theme JSON in play. Take the template's OWN accent and fonts
    where they can be read, and fall back to the readable defaults for the rest —
    so a card drawn into a corporate deck still looks like that deck."""
    theme = _apply_size_defaults(load_theme(theme_path))
    try:
        part = prs.slide_masters[0].part.part_related_by(_THEME_REL)
        xml = part.blob.decode("utf-8", "ignore")
    except Exception:
        return theme
    m = re.search(r'<a:accent1>.*?val="([0-9A-Fa-f]{6})"', xml, re.S)
    if m:
        theme["color"]["accent"] = m.group(1).upper()
    for tag, key in (("majorFont", "heading"), ("minorFont", "body")):
        f = re.search(r'<a:%s>\s*<a:latin typeface="([^"]+)"' % tag, xml)
        if f and f.group(1):
            theme["font"][key] = f.group(1)
    theme["font"].setdefault("number", theme["font"]["heading"])
    return theme


def render_template(prs, spec, map_cfg):
    slides = spec.get("slides") or []
    role_layout = _layout_index_by_role(prs)
    chosen = []
    del _WARNINGS[:]
    # Composed archetypes are drawn, not filled, so they need tokens even here.
    part_theme = _theme_of_template(prs)
    part_theme["aspect"] = ("4:3" if (prs.slide_height or 0) /
                            max(1, prs.slide_width or 1) > 0.7 else "16:9")
    part_g = make_grid(part_theme)

    for i, s in enumerate(slides, start=1):
        t = s.get("type", "bullets")
        m = (map_cfg or {}).get(t, {})
        # layout: explicit per-slide > map > heuristic
        li = s.get("layout", m.get("layout", role_layout.get(t, 0)))
        layout = _resolve_layout(prs, li)
        slide = prs.slides.add_slide(layout)
        roles = _placeholders_by_role(slide)
        chosen.append((t, layout.name))

        def pick(role, default_ph):
            if role in m:  # explicit placeholder idx from the map wins
                return _ph_by_idx(slide, m[role])
            return default_ph

        title_text = s.get("title")
        if t == "quote":
            title_text = None  # quote has no title field
        elif t == "statement":
            title_text = s.get("text") or s.get("title")
        _set_ph_text(_tpl_need(pick("title", roles["title"]), i, layout.name,
                               "title", title_text), title_text)
        if t == "quote":
            _drop_placeholder(roles["title"])   # no title text — don't leave it empty

        if t == "title":
            _set_ph_text(_tpl_need(pick("subtitle", roles["subtitle"]
                                        or (roles["body"][0] if roles["body"] else None)),
                                   i, layout.name, "subtitle", s.get("subtitle")),
                         s.get("subtitle"))
        elif t == "section":
            pass  # title placeholder already filled
        elif t == "table":
            _render_table_template(slide, s, i, layout.name)
        elif t == "chart":
            _render_chart_template(slide, s, i, layout.name)
        elif t == "two_col":
            bodies = roles["body"]
            left = _tpl_need(pick("left", bodies[0] if len(bodies) > 0 else None),
                             i, layout.name, "body", s.get("left"))
            right = pick("right", bodies[1] if len(bodies) > 1 else None)
            _fill_col_placeholder(left, s.get("left") or {})
            if right is not None:
                _fill_col_placeholder(right, s.get("right") or {})
            elif left is not None:  # no second body: merge into one
                _append_col_placeholder(left, s.get("right") or {})
        elif t == "big_number":
            body = _tpl_need(pick("body", roles["body"][0] if roles["body"] else None),
                             i, layout.name, "body", s.get("number"))
            parts = [str(s.get("number", ""))]
            if s.get("caption"):
                parts.append(s["caption"])
            _set_ph_bullets(body, parts)
            _set_ph_text(pick("source", None), s.get("source"))
        elif t == "quote":
            body = _tpl_need(pick("body", roles["body"][0] if roles["body"] else None),
                             i, layout.name, "body", s.get("quote"))
            lines = ["“" + s.get("quote", "") + "”"]
            if s.get("attribution"):
                lines.append("— " + s["attribution"])
            _set_ph_bullets(body, lines)
        elif t == "image":
            pic = pick("image", roles["picture"])
            img = _asset(s.get("image"))
            if pic is not None and img and os.path.exists(img):
                try:
                    pic.insert_picture(img)
                except Exception:
                    slide.shapes.add_picture(img, pic.left, pic.top, height=pic.height)
            elif img and os.path.exists(img):
                _warn("slide %d: layout %r has no picture placeholder — the image was placed"
                      " free-floating. Pin a layout with a PICTURE placeholder via `layout:`."
                      % (i, layout.name))
                slide.shapes.add_picture(img, Inches(1), Inches(1.5))
            _set_ph_text(pick("caption", None), s.get("caption"))
        elif t == "statement":
            if s.get("sub"):
                _set_ph_text(pick("subtitle", roles["subtitle"]
                                  or (roles["body"][0] if roles["body"] else None)),
                             s["sub"])
        elif t in ("cards", "steps", "matrix", "split", "lead"):
            # The template decides WHERE (its body placeholder's region); the
            # archetype decides WHAT gets drawn there.
            {"cards": _render_cards, "steps": _render_steps, "lead": _render_lead,
             "matrix": _render_matrix, "split": _render_split}[t](
                slide, part_theme, part_g, s, i)
            _set_ph_text(pick("source", None), s.get("source"))
        elif t == "blank":
            pass
        else:  # bullets
            _set_ph_bullets(_tpl_need(pick("body", roles["body"][0] if roles["body"] else None),
                                      i, layout.name, "body", s.get("bullets")),
                            s.get("bullets"))
            _set_ph_text(pick("source", None), s.get("source"))

    return chosen


def _fill_col_placeholder(ph, col):
    if ph is None:
        return
    items = []
    if col.get("heading"):
        items.append({"text": col["heading"], "level": 0})
    items += _normalize_bullets(col.get("bullets"))
    _set_ph_bullets(ph, items)


def _append_col_placeholder(ph, col):
    if ph is None or not (col.get("heading") or col.get("bullets")):
        return
    tf = ph.text_frame
    if col.get("heading"):
        tf.add_paragraph().text = col["heading"]
    for it in _normalize_bullets(col.get("bullets")):
        p = tf.add_paragraph()
        p.text = it["text"]
        p.level = min(it["level"], 8)


def _resolve_layout(prs, ref):
    layouts = list(prs.slide_layouts)
    if isinstance(ref, int):
        return layouts[ref] if 0 <= ref < len(layouts) else layouts[0]
    if isinstance(ref, str):
        for lo in layouts:
            if (lo.name or "").lower() == ref.lower():
                return lo
        for lo in layouts:
            if ref.lower() in (lo.name or "").lower():
                return lo
    return layouts[0]


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build(spec, out, theme_path=DEFAULT_THEME, template=None, map_path=None, base_dir=None):
    global _BASE_DIR
    _BASE_DIR = base_dir
    slides = spec.get("slides") or []
    if template:
        prs = Presentation(template)
        map_cfg = None
        if map_path:
            with open(map_path, encoding="utf-8") as f:
                map_cfg = json.load(f)
        chosen = render_template(prs, spec, map_cfg)
        _apply_notes(prs, slides)
        prs.save(out)
        print("wrote %s  (template-fill: %s)" % (out, os.path.basename(template)))
        for t, name in chosen:
            print("  %-11s -> layout %r" % (t, name))
        if _WARNINGS:
            print("  %d template warning(s) above — content may be missing from the deck."
                  % len(_WARNINGS))
    else:
        meta = spec.get("meta")
        theme = _apply_size_defaults(apply_meta(load_theme(theme_path), meta), meta)
        g = make_grid(theme)
        prs = Presentation()
        prs.slide_width = Inches(g["pageW"])
        prs.slide_height = Inches(g["pageH"])
        del _WARNINGS[:]
        render_default(prs, theme, g, slides)
        _apply_notes(prs, slides)
        prs.save(out)
        print("wrote %s  (default theme: %s)" % (out, theme.get("name", "?")))
        if _WARNINGS:
            print("  %d layout warning(s) above — consider splitting/shortening those slides."
                  % len(_WARNINGS))


def main(argv=None):
    ap = argparse.ArgumentParser(description="Clean white-based .pptx generator (python-pptx).")
    ap.add_argument("spec", help="deck spec (.yaml/.yml/.json)")
    ap.add_argument("-o", "--out", required=True, help="output .pptx path")
    ap.add_argument("--theme", default=DEFAULT_THEME, help="theme JSON (default-mode only)")
    ap.add_argument("--template", help="open a real .pptx/.potx and fill its placeholders")
    ap.add_argument("--map", dest="map_path", help="role->placeholder map JSON (template mode)")
    a = ap.parse_args(argv)
    build(load_spec(a.spec), a.out, a.theme, a.template, a.map_path,
          base_dir=os.path.dirname(os.path.abspath(a.spec)))


if __name__ == "__main__":
    main()
