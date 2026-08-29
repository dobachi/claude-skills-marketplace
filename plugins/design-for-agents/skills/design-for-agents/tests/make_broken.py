#!/usr/bin/env python3
"""検出器の陰性対照を作る。清潔なデッキに欠陥を1つずつ注入する。

各欠陥は、狙った完了条件を1つ以上発火させることを意図している。
どの欠陥がどの条件を狙っているかは INJECTED に書いてある。
"""
from __future__ import annotations

import copy
import pathlib
import struct
import sys
import zlib

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

INJECTED = {
    "13pt の本文":        ["core-type-body-floor#body-not-below-floor",
                           "core-type-one-scale#sizes-from-scale"],
    "下線・中央揃え・別書体": ["core-emphasis-ladder#no-underline-emphasis",
                           "core-layout-one-alignment#single-alignment-per-slide",
                           "core-type-sans-projection#body-is-sans",
                           "core-type-two-families-max#at-most-two-families"],
    "タイトルの重複":      ["core-a11y-reading-order#titles-unique"],
    "箇条書き9個・階層3":  ["core-density-budget#bullets-within-budget",
                           "core-density-budget#nesting-depth"],
    "影付き図形":          ["core-ornament-none#no-ornament-shapes"],
    "全幅バンド":          ["core-ornament-none#no-title-band"],
    "代替テキスト無しの画像": ["core-a11y-alt-text#alt-text-present"],
    "縦横比を崩した画像":   ["core-image-resolution#no-stretched-images",
                           "core-image-resolution#resolution-sufficient"],
    "5区分の円グラフ":     ["core-chart-pie-limit#pie-slices-within-limit"],
    "2点だけの棒グラフ":    ["core-chart-follows-question#chart-earns-place"],
    "3D グラフ":           ["core-chart-chartjunk-none#no-3d-charts"],
    "7列9行の表":          ["core-density-budget#table-within-budget"],
    "縦罫・桁不揃いの表":   ["core-table-rules-minimal#no-vertical-rules",
                           "core-table-align-numbers#numbers-right-aligned",
                           "core-table-align-numbers#decimals-consistent"],
    "薄い文字色":          ["core-a11y-contrast-body#body-contrast-4-5",
                           "core-color-four-slots#only-four-slots"],
    "逆接を含むタイトル":   ["core-density-one-message#one-claim-per-slide"],
    "書式の違うアクセント2種": ["core-emphasis-ladder#one-emphasis-per-slide"],
    "行間 1.0":            ["core-type-leading#leading-is-set"],
    "グリッド外の図形":     ["core-layout-one-grid#left-edges-aligned"],
    "プレースホルダ外の文字": ["pptx-master-placeholder#content-in-placeholders"],
    "タイトルが空":        ["pptx-master-placeholder#unique-slide-titles"],
}
# 合成できないもの（実物が要る）。rules 側は manual のまま。
NOT_INJECTABLE = {
    "pptx-diagram-smartart-none#no-smartart": "python-pptx で SmartArt を作れない",
    "pptx-master-no-stock-theme#theme-is-not-stock": "実物の Office テーマが要る",
}


def png(w: int, h: int) -> bytes:
    def chunk(tag: bytes, data: bytes) -> bytes:
        c = tag + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
    raw = b"".join(b"\x00" + b"\xcc" * (w * 3) for _ in range(h))
    return (b"\x89PNG\r\n\x1a\n"
            + chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
            + chunk(b"IDAT", zlib.compress(raw)) + chunk(b"IEND", b""))


def body_shapes(slide):
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        if getattr(sh, "is_placeholder", False) and sh.placeholder_format.idx == 0:
            continue
        if sh.top is None or sh.top >= Inches(7.0):
            continue
        if sh.text_frame.text.strip():
            yield sh


def title_of(slide):
    for sh in slide.shapes:
        if getattr(sh, "is_placeholder", False) and sh.placeholder_format.idx == 0:
            return sh
    return None


def build(src: pathlib.Path, out: pathlib.Path) -> None:
    prs = Presentation(str(src))
    sl = prs.slides
    tmp = out.parent

    for sh in body_shapes(sl[1]):
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)
        break

    for sh in body_shapes(sl[3]):
        p = sh.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.0
        for r in p.runs:
            r.font.underline = True
            r.font.name = "Comic Sans MS"
        break

    a, b = title_of(sl[5]), title_of(sl[6])
    if a is not None and b is not None:
        b.text_frame.text = a.text_frame.text

    for sh in body_shapes(sl[2]):
        base = sh.text_frame.paragraphs[0]._p
        for _ in range(8):
            sh.text_frame._txBody.append(copy.deepcopy(base))
        extra = sh.text_frame.add_paragraph()
        extra.text = "第3階層"
        extra.level = 2
        break

    box = sl[8].shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1.37), Inches(2), Inches(3), Inches(1))
    eff = etree.SubElement(box._element.spPr, f"{A}effectLst")
    etree.SubElement(eff, f"{A}outerShdw", blurRad="50800", dist="38100", dir="2700000")

    band = sl[9].shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0.3), prs.slide_width, Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0x2F, 0x5D, 0xA8)
    band.line.fill.background()

    p = tmp / "_tiny.png"
    p.write_bytes(png(8, 8))
    sl[1].shapes.add_picture(str(p), Inches(9.5), Inches(1.2), Inches(3.0), Inches(1.0))

    cd = CategoryChartData()
    cd.categories = ["A", "B", "C", "D", "E"]
    cd.add_series("s", (3, 2, 2, 1, 1))
    sl[3].shapes.add_chart(XL_CHART_TYPE.PIE, Inches(9.2), Inches(3.0),
                           Inches(3.5), Inches(3.0), cd)

    cd2 = CategoryChartData()
    cd2.categories = ["X", "Y"]
    cd2.add_series("s", (1, 2))
    gf = sl[4].shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(9.2), Inches(1.2),
                                Inches(3.5), Inches(2.5), cd2)
    cs = gf.chart._chartSpace
    chart_el = cs.find(f"{{{cs.nsmap['c']}}}chart")
    etree.SubElement(chart_el, f"{{{cs.nsmap['c']}}}view3D")

    tbl = sl[5].shapes.add_table(9, 7, Inches(0.92), Inches(2.0),
                                 Inches(11.4), Inches(3.5)).table
    tbl.cell(1, 1).text = "1.5"
    tbl.cell(2, 1).text = "2.25"
    for r in (1, 2):
        for para in tbl.cell(r, 1).text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
    for cell in (tbl.cell(1, 1), tbl.cell(2, 1)):
        tcPr = cell._tc.get_or_add_tcPr()
        etree.SubElement(tcPr, f"{A}lnL", w="12700")

    for sh in body_shapes(sl[6]):
        r = sh.text_frame.paragraphs[0].runs[0]
        r.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        break

    t = title_of(sl[7])
    if t is not None:
        t.text_frame.text = "コストは下がる、しかし移行の負荷は上がる"

    for sh in body_shapes(sl[8]):
        paras = [p for p in sh.text_frame.paragraphs if p.runs]
        if len(paras) >= 2:
            paras[0].runs[0].font.color.rgb = RGBColor(0x2F, 0x5D, 0xA8)
            paras[0].runs[0].font.size = Pt(18)
            paras[1].runs[0].font.color.rgb = RGBColor(0x2F, 0x5D, 0xA8)
            paras[1].runs[0].font.size = Pt(16)
            paras[1].runs[0].font.bold = True
        break

    stray = sl[7].shapes.add_textbox(Inches(3.33), Inches(5.0), Inches(4), Inches(0.6))
    stray.text_frame.text = "プレースホルダ外の自由なテキストボックス"

    t9 = title_of(sl[9])
    if t9 is not None:
        t9.text_frame.text = ""

    prs.save(str(out))


def main(argv=None) -> int:
    argv = argv or sys.argv[1:]
    if len(argv) != 2:
        print("usage: make_broken.py <clean.pptx> <broken.pptx>", file=sys.stderr)
        return 2
    build(pathlib.Path(argv[0]), pathlib.Path(argv[1]))
    print(f"注入した欠陥 {len(INJECTED)} 種 -> {argv[1]}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
