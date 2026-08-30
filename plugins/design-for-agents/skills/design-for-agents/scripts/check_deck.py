#!/usr/bin/env python3
"""既存の .pptx を、rules/ の完了条件に照らして検査する。

各検査の id は `<rule-id>#<check-id>` であり、rules/ の `done_when` と
1対1に対応する。閾値は tokens/pptx.tokens.json から読む（値をここに書かない）。

検証の状態:
  正常系（pptx-build のサンプル3本）で偽陽性ゼロ、欠陥を注入した異常系で全て
  発火することを確認した条件だけ、rules/ 側を `check: automated` にしてある。
  構図の3件（同一骨格の連続・地の暗いページ・角丸と線幅）は 2026-08-30 に追加した。
  以下2件は実装済みだが**発火を確認していない**ので、rules/ 側は manual のまま。
    - pptx-diagram-smartart-none#no-smartart（python-pptx で SmartArt を作れない）
    - pptx-master-no-stock-theme#theme-is-not-stock（実物の Office テーマが要る）

終了コード:
  0  違反なし
  1  違反あり
  2  検査できない（ファイルが読めない、下限を割っている）

使い方:
  python3 scripts/check_deck.py deck.pptx
  python3 scripts/check_deck.py deck.pptx --json
"""
from __future__ import annotations

import argparse
import colorsys
import json
import pathlib
import sys

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:  # pragma: no cover
    print("python-pptx が要る: pip install -r assets/requirements.txt", file=sys.stderr)
    sys.exit(2)

ROOT = pathlib.Path(__file__).resolve().parent.parent
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# 出典行は本文ではない（core-source-line-on-data / core-type-body-floor の
# not_applicable_when）。フッタ帯の外（図の直下など）にも置かれるので、
# 位置ではなく書き出しで識別する。
SOURCE_PREFIXES = ("出典", "出所", "Source", "source", "Note", "注:", "注：")

# 下限: 対象を消して「違反なし」にするのを防ぐ
FLOOR_BODY_SLIDES = 3

ORNAMENT_TAGS = ("gradFill", "outerShdw", "innerShdw", "bevelT", "bevelB", "glow", "reflection")
STOCK_THEMES = {"ion", "facet", "wisp", "berlin", "celestial", "damask", "depth",
                "dividend", "droplet", "frame", "mesh", "metropolitan", "parallax",
                "quotable", "retrospect", "savon", "slice", "vapor trail", "view"}


def load_tokens() -> dict:
    return json.loads((ROOT / "tokens" / "pptx.tokens.json").read_text(encoding="utf-8"))


class Findings:
    def __init__(self) -> None:
        self.rows: list[tuple[str, int, str]] = []   # (check_id, slide_no, message)
        self.skipped: list[tuple[str, str]] = []     # (check_id, なぜ評価しなかったか)

    def add(self, check: str, slide: int, message: str) -> None:
        self.rows.append((check, slide, message))

    def skip(self, check: str, why: str) -> None:
        self.skipped.append((check, why))

    def of(self, check: str) -> list:
        return [r for r in self.rows if r[0] == check]


# ---------------------------------------------------------------- 走査の下ごしらえ

def iter_text(shape):
    """(段落, run) を返す。図形が文字を持たなければ何も返さない。"""
    if not getattr(shape, "has_text_frame", False):
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            yield para, run


def shapes_of(slide):
    """グループを展開して全図形を返す。"""
    out = []
    def walk(container):
        for sh in container:
            if sh.shape_type == 6 and hasattr(sh, "shapes"):   # GROUP
                walk(sh.shapes)
            else:
                out.append(sh)
    walk(slide.shapes)
    return out


def is_smartart(shape) -> bool:
    xml = shape._element.xml
    return "graphicData" in xml and "diagramLayout" in xml


def alt_text_of(shape) -> str:
    el = shape._element
    nv = el.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
    if nv is None:
        nv = el.find(f".//{NS_A}cNvPr")
    return (nv.get("descr") or "").strip() if nv is not None else ""


def in_footer(shape, tok) -> bool:
    """フッタ線より下にある要素か。出典行とページ番号がここに入る。

    rules の not_applicable_when が「出典行・ページ番号は本文ではない」と
    定めているので、本文向けの検査から外す。
    """
    top = getattr(shape, "top", None)
    if top is None:
        return False
    return top >= Emu(int(tok["grid"]["footerTop"]["$value"] * 914400))


def is_source_line(text: str) -> bool:
    t = text.strip()
    return any(t.startswith(p) for p in SOURCE_PREFIXES)


def is_body_text(shape, tok) -> bool:
    """本文とみなす図形か。タイトル・フッタ・表の中・出典行は本文ではない。"""
    if not getattr(shape, "has_text_frame", False):
        return False
    if getattr(shape, "has_table", False):
        return False
    if getattr(shape, "is_placeholder", False) and shape.placeholder_format.idx == 0:
        return False
    if is_source_line(shape.text_frame.text):
        return False
    return not in_footer(shape, tok)


def body_slide_count(prs) -> int:
    """タイトルと本文の両方を持つスライドの数（下限の判定に使う）。"""
    n = 0
    for slide in prs.slides:
        has_title = any(sh.has_text_frame and sh.text_frame.text.strip()
                        and sh.is_placeholder and sh.placeholder_format.idx == 0
                        for sh in slide.shapes if hasattr(sh, "is_placeholder"))
        has_body = any(sh.has_text_frame and sh.text_frame.text.strip()
                       for sh in slide.shapes
                       if getattr(sh, "has_text_frame", False)
                       and not (getattr(sh, "is_placeholder", False)
                                and sh.placeholder_format.idx == 0))
        if has_title and has_body:
            n += 1
    return n


# ---------------------------------------------------------------- 検査

def check_type_scale(prs, tok, f):
    allowed = {v["$value"] for k, v in tok["fontSize"].items() if k != "$type"
               and isinstance(v, dict) and "$value" in v}
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            for _, run in iter_text(sh):
                if run.font.size is None:      # レイアウト継承。マスターが決めている
                    continue
                pt = round(run.font.size.pt, 1)
                if pt not in allowed:
                    f.add("core-type-one-scale#sizes-from-scale", i,
                          f"型スケールに無いサイズ {pt}pt「{run.text[:18]}」")


def check_body_floor(prs, tok, f):
    floor = tok["fontSize"]["body"]["$value"]
    sub = tok["fontSize"]["bodySub"]["$value"]
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if not is_body_text(sh, tok):       # タイトル・フッタ・表は本文ではない
                continue
            for para, run in iter_text(sh):
                if run.font.size is None:
                    continue
                pt = run.font.size.pt
                limit = sub if para.level >= 1 else floor
                if pt < limit:
                    f.add("core-type-body-floor#body-not-below-floor", i,
                          f"本文 {pt}pt < 下限 {limit}pt「{run.text[:18]}」")


def check_font_families(prs, tok, f):
    fams = set()
    for slide in prs.slides:
        for sh in shapes_of(slide):
            for _, run in iter_text(sh):
                if run.font.name:
                    fams.add(run.font.name)
    if len(fams) > 2:
        f.add("core-type-two-families-max#at-most-two-families", 0,
              f"書体ファミリーが {len(fams)} 種: {', '.join(sorted(fams))}")


def check_density(prs, tok, f):
    top_max = tok["density"]["bulletsTopLevelMax"]["$value"]
    depth_max = tok["density"]["bulletDepthMax"]["$value"]
    cols_max = tok["density"]["tableColumnsMax"]["$value"]
    rows_max = tok["density"]["tableRowsMax"]["$value"]
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if getattr(sh, "has_table", False):
                t = sh.table
                if len(t.columns) > cols_max or len(t.rows) > rows_max:
                    f.add("core-density-budget#table-within-budget", i,
                          f"表 {len(t.columns)}列 × {len(t.rows)}行 が上限 "
                          f"{cols_max}×{rows_max} を超える")
            if not is_body_text(sh, tok):
                continue
            paras = [p for p in sh.text_frame.paragraphs if p.text.strip()]
            top = sum(1 for p in paras if p.level == 0)
            if top > top_max:
                f.add("core-density-budget#bullets-within-budget", i,
                      f"第1階層の箇条書きが {top} 個（上限 {top_max}）")
            deepest = max((p.level for p in paras), default=0) + 1
            if deepest > depth_max:
                f.add("core-density-budget#nesting-depth", i,
                      f"箇条書きの階層が {deepest}（上限 {depth_max}）")


def check_ornament(prs, tok, f):
    slide_w = prs.slide_width
    title_bottom = Emu(int(tok["grid"]["titleTop"]["$value"] * 914400)) + \
        Emu(int(tok["grid"]["titleHeight"]["$value"] * 914400))
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            xml = sh._element.xml
            hit = [t for t in ORNAMENT_TAGS if f"<{NS_A[1:-1]}" not in xml and f"a:{t}" in xml]
            if hit:
                f.add("core-ornament-none#no-ornament-shapes", i,
                      f"装飾効果 {', '.join(sorted(set(hit)))} を持つ図形")
            # 自動図形は空のテキストフレームを持つので、has_text_frame では外せない。
            # 文字が入っているものだけを本文として除外する。
            has_words = (getattr(sh, "has_text_frame", False)
                         and sh.text_frame.text.strip())
            if has_words or getattr(sh, "has_table", False):
                continue
            if sh.width is None or sh.top is None:
                continue
            if sh.width >= slide_w * 0.95 and sh.top < title_bottom and "solidFill" in xml:
                f.add("core-ornament-none#no-title-band", i,
                      "タイトル領域にスライド幅いっぱいの塗り図形")


def check_underline(prs, tok, f):
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            for _, run in iter_text(sh):
                if run.font.underline:
                    f.add("core-emphasis-ladder#no-underline-emphasis", i,
                          f"下線「{run.text[:18]}」")


def check_alignment(prs, tok, f):
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if not is_body_text(sh, tok):
                continue
            for para in sh.text_frame.paragraphs:
                if not para.text.strip() or para.alignment is None:
                    continue
                if str(para.alignment).startswith(("CENTER", "RIGHT", "JUSTIFY")):
                    f.add("core-layout-one-alignment#single-alignment-per-slide", i,
                          f"本文が {str(para.alignment).split()[0]} 揃え「{para.text[:18]}」")


def check_pie(prs, tok, f):
    limit = tok["density"]["pieSlicesMax"]["$value"]
    seen = False
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if not getattr(sh, "has_chart", False):
                continue
            seen = True
            if "pie" not in str(sh.chart.chart_type).lower():
                continue
            n = len(list(sh.chart.plots[0].categories))
            if n > limit:
                f.add("core-chart-pie-limit#pie-slices-within-limit", i,
                      f"円グラフの区分が {n}（上限 {limit}）")
    if not seen:
        f.skip("core-chart-pie-limit#pie-slices-within-limit", "グラフが1つも無い")


def check_titles_unique(prs, tok, f):
    seen: dict[str, int] = {}
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if getattr(sh, "is_placeholder", False) and sh.placeholder_format.idx == 0:
                t = sh.text_frame.text.strip()
                if not t:
                    continue
                if t in seen:
                    f.add("core-a11y-reading-order#titles-unique", i,
                          f"タイトルが {seen[t]} 枚目と同一「{t[:24]}」")
                seen.setdefault(t, i)


def check_smartart(prs, tok, f):
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if is_smartart(sh):
                f.add("pptx-diagram-smartart-none#no-smartart", i, "SmartArt 図形")


def check_alt_text(prs, tok, f):
    seen = False
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if getattr(sh, "has_text_frame", False) or getattr(sh, "has_table", False):
                continue
            if sh.shape_type is None:
                continue
            seen = True
            if not alt_text_of(sh):
                f.add("core-a11y-alt-text#alt-text-present", i,
                      f"代替テキストが空の要素（{sh.shape_type}）")
    if not seen:
        f.skip("core-a11y-alt-text#alt-text-present", "文字以外の要素が1つも無い")


def check_stock_theme(prs, tok, f):
    for master in prs.slide_masters:
        part = master.part
        for rel in part.rels.values():
            if "theme" in rel.reltype:
                name = ""
                try:
                    name = rel.target_part.blob.decode("utf-8", "ignore")[:400]
                except Exception:
                    continue
                for stock in STOCK_THEMES:
                    if f'name="{stock}"' in name.lower():
                        f.add("pptx-master-no-stock-theme#theme-is-not-stock", 0,
                              f"同梱テーマ「{stock}」が設定されている")



# ---------------------------------------------------------------- 第2バッチ

def _rgb(color):
    """明示された RGB だけを返す。テーマ色や未指定は None。"""
    try:
        if color is None or color.type is None:
            return None
        if str(color.type) != "MSO_THEME_COLOR.NOT_THEME_COLOR" and "SCHEME" in str(color.type):
            return None
        return str(color.rgb).upper()
    except Exception:
        return None


def _relative_luminance(hex6: str) -> float:
    c = [int(hex6[i:i + 2], 16) / 255 for i in (0, 2, 4)]
    c = [x / 12.92 if x <= 0.04045 else ((x + 0.055) / 1.055) ** 2.4 for x in c]
    return 0.2126 * c[0] + 0.7152 * c[1] + 0.0722 * c[2]


def contrast_ratio(fg: str, bg: str) -> float:
    l1, l2 = sorted((_relative_luminance(fg), _relative_luminance(bg)), reverse=True)
    return (l1 + 0.05) / (l2 + 0.05)


def deck_background(prs, tok) -> str:
    """スライドの地色。明示されていなければトークンの paper を使う。"""
    return tok["color"]["paper"]["$value"].lstrip("#").upper()


def check_leading(prs, tok, f):
    want = tok["text"]["lineHeight"]["$value"]
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if not is_body_text(sh, tok):
                continue
            for para in sh.text_frame.paragraphs:
                ls = para.line_spacing
                if ls is None or not isinstance(ls, float):
                    continue                     # 継承。マスターが決めている
                if abs(ls - want) > 0.01:
                    f.add("core-type-leading#leading-is-set", i,
                          f"行間 {ls} が定めた {want} と違う「{para.text[:18]}」")


def check_sans(prs, tok, f):
    allowed = set(tok["fontFamily"]["sans"]["$value"])
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if not is_body_text(sh, tok):
                continue
            for _, run in iter_text(sh):
                name = run.font.name
                if name and name not in allowed:
                    f.add("core-type-sans-projection#body-is-sans", i,
                          f"本文の書体『{name}』が宣言したサンセリフに無い")


def check_grid(prs, tok, f):
    g = tok["grid"]
    margin = g["sideMargin"]["$value"]
    content = g["contentWidth"]["$value"]
    gutter = g["gutter"]["$value"]
    col2 = margin + (content - gutter) / 2 + gutter
    allowed = [0.0, margin, round(col2, 2)]
    tol = 0.02
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if sh.left is None:
                continue
            left = round(sh.left / 914400, 2)
            if any(abs(left - a) <= tol for a in allowed):
                continue
            # 右揃えの要素は右端で揃う（ページ番号など）
            right_margin = round(prs.slide_width / 914400 - margin, 2)
            right = round((sh.left + (sh.width or 0)) / 914400, 2)
            if abs(right - right_margin) <= tol:
                continue
            f.add("core-layout-one-grid#left-edges-aligned", i,
                  f"左端 {left}in も右端 {right}in もグリッドに乗っていない")


def check_four_slots(prs, tok, f):
    """色数を数える。色値は見ない — ブランドがあれば accent は差し替わるため。"""
    limit = 3                                     # ink / muted / accent。紙は背景
    used: dict[str, int] = {}
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if getattr(sh, "has_table", False) or getattr(sh, "has_chart", False):
                continue                          # 表とグラフの内部は対象外
            for _, run in iter_text(sh):
                rgb = _rgb(run.font.color)
                if rgb:
                    used.setdefault(rgb, i)
    if len(used) > limit:
        f.add("core-color-four-slots#only-four-slots", 0,
              f"文字色が {len(used)} 種類（上限 {limit}）: "
              + ", ".join(f"#{c}" for c in sorted(used)))


def check_one_emphasis(prs, tok, f):
    """書式が違うアクセント色が同一スライドに2種類あれば強調の競合。

    同じ大きさ・太さで並ぶもの（2段組の両列見出しなど）は反復であって
    競合ではないので、1箇所と数える（core-emphasis-ladder の
    not_applicable_when）。デッキ内で最も多い色を accent とみなすのではなく、
    トークンの accent と、デッキが上書きした accent の両方を拾う。
    """
    ink = tok["color"]["ink"]["$value"].lstrip("#").upper()
    muted = tok["color"]["muted"]["$value"].lstrip("#").upper()
    for i, slide in enumerate(prs.slides, 1):
        sigs = set()
        for sh in shapes_of(slide):
            if getattr(sh, "has_table", False) or getattr(sh, "has_chart", False):
                continue
            for _, run in iter_text(sh):
                rgb = _rgb(run.font.color)
                if not rgb or rgb in (ink, muted) or not run.text.strip():
                    continue
                pt = run.font.size.pt if run.font.size else None
                sigs.add((rgb, pt, bool(run.font.bold)))
        if len(sigs) > 1:
            f.add("core-emphasis-ladder#one-emphasis-per-slide", i,
                  f"書式の違うアクセントが {len(sigs)} 種類: "
                  + ", ".join(f"#{c}/{s}pt/{'B' if b else '-'}" for c, s, b in sorted(sigs, key=str)))


def check_contrast(prs, tok, f):
    bg = deck_background(prs, tok)
    large = 18.0
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if getattr(sh, "has_chart", False):
                continue
            for _, run in iter_text(sh):
                fg = _rgb(run.font.color)
                if not fg or not run.text.strip():
                    continue
                pt = run.font.size.pt if run.font.size else large
                need = 3.0 if (pt >= large or (pt >= 14 and run.font.bold)) else 4.5
                r = contrast_ratio(fg, bg)
                if r < need:
                    f.add("core-a11y-contrast-body#body-contrast-4-5", i,
                          f"コントラスト比 {r:.2f} < {need}（#{fg} on #{bg}, {pt}pt）")


def check_tables(prs, tok, f):
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if not getattr(sh, "has_table", False):
                continue
            tbl = sh.table
            if "lnL" in sh._element.xml or "lnR" in sh._element.xml:
                f.add("core-table-rules-minimal#no-vertical-rules", i, "表に縦罫がある")
            ncol = len(tbl.columns)
            for c in range(ncol):
                cells = [tbl.cell(r, c) for r in range(1, len(tbl.rows))]
                nums = [(x, x.text.strip()) for x in cells
                        if _is_number(x.text.strip())]
                if not nums:
                    continue
                for cell, _t in nums:
                    for para in cell.text_frame.paragraphs:
                        if para.text.strip() and para.alignment is not None \
                           and not str(para.alignment).startswith("RIGHT"):
                            f.add("core-table-align-numbers#numbers-right-aligned", i,
                                  f"{c + 1}列目の数値が右揃えでない「{para.text[:12]}」")
                decs = {len(t.split(".")[1]) if "." in t else 0 for _c, t in nums}
                if len(decs) > 1:
                    f.add("core-table-align-numbers#decimals-consistent", i,
                          f"{c + 1}列目の小数点以下の桁数が揃っていない {sorted(decs)}")


def _is_number(t: str) -> bool:
    t = t.replace(",", "").replace("%", "").replace("円", "").strip()
    if not t:
        return False
    try:
        float(t)
        return True
    except ValueError:
        return False


def check_charts(prs, tok, f):
    seen = False
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if not getattr(sh, "has_chart", False):
                continue
            seen = True
            xml = sh.chart._chartSpace.xml
            if "view3D" in xml:
                f.add("core-chart-chartjunk-none#no-3d-charts", i, "3D 効果を持つグラフ")
            if xml.count("<c:valAx>") > 1:
                f.add("core-chart-chartjunk-none#no-dual-axis", i, "第2軸を持つグラフ")
            try:
                n = len(list(sh.chart.plots[0].categories))
            except Exception:
                n = 0
            if 0 < n < 3:
                f.add("core-chart-follows-question#chart-earns-place", i,
                      f"データ点が {n} 個のグラフ（3個未満はグラフにしない）")
    if not seen:
        for c in ("core-chart-chartjunk-none#no-3d-charts",
                  "core-chart-chartjunk-none#no-dual-axis",
                  "core-chart-follows-question#chart-earns-place"):
            f.skip(c, "グラフが1つも無い")


def check_images(prs, tok, f):
    seen = False
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            if not hasattr(sh, "image"):
                continue
            seen = True
            try:
                px_w, px_h = sh.image.size
            except Exception:
                continue
            if not sh.width or not sh.height:
                continue
            native = px_w / px_h
            shown = sh.width / sh.height
            if abs(native - shown) / native > 0.02:
                f.add("core-image-resolution#no-stretched-images", i,
                      f"縦横比が元画像と違う（元 {native:.2f} / 表示 {shown:.2f}）")
            dpi = px_w / (sh.width / 914400)
            if dpi < 150:
                f.add("core-image-resolution#resolution-sufficient", i,
                      f"表示サイズあたり {dpi:.0f}dpi（150 未満）")
    if not seen:
        for c in ("core-image-resolution#no-stretched-images",
                  "core-image-resolution#resolution-sufficient"):
            f.skip(c, "画像が1つも無い")


def check_placeholders(prs, tok, f):
    for i, slide in enumerate(prs.slides, 1):
        stray = 0
        has_title = False
        for sh in slide.shapes:
            if getattr(sh, "is_placeholder", False):
                if sh.placeholder_format.idx == 0 and sh.has_text_frame \
                   and sh.text_frame.text.strip():
                    has_title = True
                continue
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip() \
               and not in_footer(sh, tok) and not is_source_line(sh.text_frame.text):
                stray += 1
        if stray:
            f.add("pptx-master-placeholder#content-in-placeholders", i,
                  f"プレースホルダ外の文字要素が {stray} 個")
        if not has_title:
            f.add("pptx-master-placeholder#unique-slide-titles", i,
                  "タイトルプレースホルダが空")


CONNECTIVES = ("しかし", "一方で", "だが", "および", "かつ", "ならびに")


def check_one_claim(prs, tok, f):
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if getattr(sh, "is_placeholder", False) and sh.placeholder_format.idx == 0:
                t = sh.text_frame.text.strip()
                hit = [c for c in CONNECTIVES if c in t]
                if hit:
                    f.add("core-density-one-message#one-claim-per-slide", i,
                          f"タイトルが「{hit[0]}」で2つの主張を繋いでいる「{t[:28]}」")


# ---------------------------------------------------------------- 構図と造形

def _hls(hex6: str):
    hex6 = str(hex6).lstrip("#")
    r, g, b = (int(hex6[i:i + 2], 16) / 255.0 for i in (0, 2, 4))
    return colorsys.rgb_to_hls(r, g, b)


def _ground_lightness(slide) -> float:
    """スライドが自分で塗った地の明度。塗っていなければ白とみなす。"""
    try:
        rgb = slide.background.fill.fore_color.rgb
        if rgb is not None:
            return _hls("%02X%02X%02X" % (rgb[0], rgb[1], rgb[2]))[1]
    except Exception:
        pass
    return 1.0


def _skeleton(slide) -> tuple:
    """遠目に見たときのスライドの姿。レイアウト名、載っている部品の種類、
    図表の有無。同じ姿のスライドは同じスライドに見える。"""
    kinds = tuple(sorted({(sh.name or "").split("/", 1)[1]
                          for sh in slide.shapes
                          if (sh.name or "").startswith("part/")}))
    gfx = any(getattr(sh, "has_table", False) or getattr(sh, "has_chart", False)
              for sh in slide.shapes)
    pic = any("PICTURE" in str(sh.shape_type) for sh in slide.shapes)
    return (slide.slide_layout.name, kinds, gfx, pic)


def check_skeleton_run(prs, tok, f):
    """同じ骨格が続くと、内容の違いをレイアウトが均してしまう。"""
    limit = int(tok["archetype"]["sameSkeletonRunMax"]["$value"])
    if body_slide_count(prs) < 6:
        f.skip("core-structure-skeleton-varies#no-long-identical-run",
               "本文スライドが6枚未満（下限を割っている）")
        return
    skel = [_skeleton(s) for s in prs.slides]
    run, start = 1, 1
    for i in range(1, len(skel) + 1):
        same = i < len(skel) and skel[i] == skel[i - 1]
        if same:
            run += 1
            continue
        if run > limit:
            f.add("core-structure-skeleton-varies#no-long-identical-run", start,
                  f"同じ骨格が {run} 枚続く（{skel[start - 1][0]}）")
        start, run = i + 1, 1


def check_dark_pages(prs, tok, f):
    """転換の印は、多用すれば印ではなくなる。"""
    limit = int(tok["archetype"]["darkPagesMax"]["$value"])
    if body_slide_count(prs) < 6:
        f.skip("core-emphasis-dark-page#dark-pages-limited",
               "本文スライドが6枚未満（下限を割っている）")
        return
    grounds = [_ground_lightness(s) for s in prs.slides]
    common = max(set(grounds), key=grounds.count)      # デッキの地
    odd = [i for i, g in enumerate(grounds, 1) if abs(g - common) > 0.25]
    if len(odd) > limit:
        f.add("core-emphasis-dark-page#dark-pages-limited", odd[0],
              f"地の異なるページが {len(odd)} 枚（上限 {limit}）")


def _corner_radius(shape):
    """角丸の調整値。角丸を持たない図形は None。"""
    try:
        adj = shape.adjustments
        if len(adj) == 0:
            return None
        return round(float(adj[0]), 4)
    except Exception:
        return None


def check_shape_tokens(prs, tok, f):
    """角丸と線幅がデッキで1種類か。2種類あると、理由は言えないのに雑に見える。"""
    radii, widths = {}, {}
    for i, slide in enumerate(prs.slides, 1):
        for sh in shapes_of(slide):
            r = _corner_radius(sh)
            if r is not None and r > 0:
                radii.setdefault(r, i)
            try:
                w = sh.line.width
            except Exception:
                w = None
            if w:
                widths.setdefault(round(Emu(int(w)).pt, 2), i)
    if len(radii) < 2 and len(widths) < 2:
        if not radii and not widths:
            f.skip("pptx-ornament-three-tokens#one-corner-radius",
                   "角丸または線を持つ図形が無い（下限を割っている）")
    if len(radii) > 1:
        f.add("pptx-ornament-three-tokens#one-corner-radius", min(radii.values()),
              "角丸が %d 種類ある（%s）" % (len(radii), ", ".join(str(k) for k in sorted(radii))))
    if len(widths) > 1:
        f.add("pptx-ornament-three-tokens#one-line-weight", min(widths.values()),
              "線幅が %d 種類ある（%s pt）" % (len(widths), ", ".join(str(k) for k in sorted(widths))))


CHECKS = [
    ("core-type-one-scale#sizes-from-scale", check_type_scale),
    ("core-type-body-floor#body-not-below-floor", check_body_floor),
    ("core-type-two-families-max#at-most-two-families", check_font_families),
    ("core-density-budget#bullets-within-budget", check_density),
    ("core-ornament-none#no-ornament-shapes", check_ornament),
    ("core-emphasis-ladder#no-underline-emphasis", check_underline),
    ("core-layout-one-alignment#single-alignment-per-slide", check_alignment),
    ("core-chart-pie-limit#pie-slices-within-limit", check_pie),
    ("core-a11y-reading-order#titles-unique", check_titles_unique),
    ("pptx-diagram-smartart-none#no-smartart", check_smartart),
    ("core-a11y-alt-text#alt-text-present", check_alt_text),
    ("pptx-master-no-stock-theme#theme-is-not-stock", check_stock_theme),
    ("core-type-leading#leading-is-set", check_leading),
    ("core-type-sans-projection#body-is-sans", check_sans),
    ("core-layout-one-grid#left-edges-aligned", check_grid),
    ("core-color-four-slots#only-four-slots", check_four_slots),
    ("core-emphasis-ladder#one-emphasis-per-slide", check_one_emphasis),
    ("core-a11y-contrast-body#body-contrast-4-5", check_contrast),
    ("core-table-rules-minimal#no-vertical-rules", check_tables),
    ("core-chart-chartjunk-none#no-3d-charts", check_charts),
    ("core-image-resolution#no-stretched-images", check_images),
    ("pptx-master-placeholder#content-in-placeholders", check_placeholders),
    ("core-density-one-message#one-claim-per-slide", check_one_claim),
    ("core-structure-skeleton-varies#no-long-identical-run", check_skeleton_run),
    ("core-emphasis-dark-page#dark-pages-limited", check_dark_pages),
    ("pptx-ornament-three-tokens#one-corner-radius", check_shape_tokens),
]


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description="rules/ の完了条件で .pptx を検査する")
    ap.add_argument("deck")
    ap.add_argument("--json", action="store_true", help="機械可読で出す")
    args = ap.parse_args(argv)

    path = pathlib.Path(args.deck)
    if not path.is_file():
        print(f"開けない: {path}", file=sys.stderr)
        return 2
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        print(f".pptx として読めない: {exc}", file=sys.stderr)
        return 2

    tok = load_tokens()
    f = Findings()

    # 下限を先に見る。対象を消して「違反なし」にできないようにする
    bodies = body_slide_count(prs)
    if bodies < FLOOR_BODY_SLIDES:
        msg = (f"下限を割っている: タイトルと本文を持つスライドが {bodies} 枚 "
               f"（{FLOOR_BODY_SLIDES} 枚以上が必要）")
        print(json.dumps({"exit": 2, "reason": msg}, ensure_ascii=False) if args.json else msg,
              file=sys.stderr)
        return 2

    seen = set()
    for _, fn in CHECKS:
        if fn in seen:
            continue
        seen.add(fn)
        fn(prs, tok, f)

    if args.json:
        print(json.dumps({
            "deck": str(path), "slides": len(prs.slides), "body_slides": bodies,
            "violations": [{"check": c, "slide": s, "message": m} for c, s, m in f.rows],
            "skipped": [{"check": c, "why": w} for c, w in f.skipped],
        }, ensure_ascii=False, indent=2))
    else:
        print(f"{path.name}: {len(prs.slides)} スライド（本文 {bodies} 枚）"
              f" / 検査 {len(CHECKS)} 件")
        for c, w in f.skipped:
            print(f"  対象なし {c} — {w}")
        for c, s, m in f.rows:
            where = f"スライド{s}" if s else "デッキ全体"
            print(f"  違反 {c}  [{where}] {m}", file=sys.stderr)
        print("違反なし" if not f.rows else f"\n違反 {len(f.rows)} 件", 
              file=sys.stderr if f.rows else sys.stdout)
    return 1 if f.rows else 0


if __name__ == "__main__":
    sys.exit(main())
